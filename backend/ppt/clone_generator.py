"""
Clone-based Slide Generator for LivePPT.

Implements PPTAgent-inspired clone-and-edit approach:
1. Deep copy a template slide (preserving ALL visual formatting)
2. Replace text content paragraph-by-paragraph
3. Handle paragraph clone/delete for variable-length elements
4. Preserve decorations, images, groups, auto-shapes untouched

Works with text-only LLM — no vision model required.
Paired with clone_inducter.py which provides TemplateSpec.
"""

import json
import re
from copy import deepcopy
from typing import List, Dict, Optional, Any, Tuple
from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.llm_config import LLMConfig
from backend.ppt.clone_inducter import (
    TemplateInducter,
    TemplateSpec,
    InducedLayout,
    SlideCategory,
    ContentSchema,
    TextElementInfo,
)


# ═══════════════════════════════════════════════════════════════
#  Content Filler: LLM generates content for a schema
# ═══════════════════════════════════════════════════════════════

CONTENT_FILL_SYSTEM = """You are a world-class presentation content writer.
Given a slide layout schema and a topic, generate compelling, insight-rich text for each element.

Your writing philosophy:
- Every bullet point must contain SPECIFIC information: who, what, how much, why it matters
- Titles should be CONCLUSIONS, not topic labels ("Supply chain attacks up 300%, legacy defenses failing" NOT "Security threats")
- Use contrast and data to create cognitive impact
- Build logical flow between points: cause→effect, problem→solution, old→new
- NEVER write vague platitudes like "strengthen security" or "drive innovation"

Rules:
- Respect the suggested_characters for each element — match the length as closely as possible
- For variable_length elements, provide the appropriate number of items (usually 3-5)
- Output valid JSON matching the schema
- Write in the same language as the input content
- Extract real data and facts from the source content — never fabricate statistics
- Each text element should deliver a distinct piece of information, not repeat the same point"""

CONTENT_FILL_PROMPT = """## Slide Layout Schema
{schema}

## Slide Purpose
{purpose}

## Source Content
{content}

## Instructions
Generate HIGH-QUALITY content for each element. Remember:
- Titles = core conclusions or insights, not generic labels
- Bullets = specific facts with data/evidence, not vague statements
- Each element should advance the audience's understanding
- Use the source content's real data, examples, and insights

Output JSON:
```json
{{
  "element_name": ["text1", "text2", ...],
  ...
}}
```

For single-value elements, still use a list with one item.
For image elements, output an empty list."""


class ContentFiller:
    """Use LLM to generate content matching a layout schema."""

    def fill(self, schema: ContentSchema, purpose: str, content: str) -> Dict[str, List[str]]:
        """
        Generate content for each schema element.
        
        Args:
            schema: The layout's content schema
            purpose: What this slide is about
            content: Source material to draw from
            
        Returns:
            Dict mapping element name to list of strings
        """
        config = LLMConfig()
        if not config.is_configured():
            raise RuntimeError("LLM 未配置")

        prompt = CONTENT_FILL_PROMPT.format(
            schema=schema.to_prompt(),
            purpose=purpose,
            content=content[:3000],  # Limit content length
        )

        import requests as http_requests
        url = f"{config.config['base_url'].rstrip('/')}/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {config.config['api_key']}",
        }
        payload = {
            "model": config.config["model"],
            "messages": [
                {"role": "system", "content": CONTENT_FILL_SYSTEM},
                {"role": "user", "content": prompt},
            ],
            "max_tokens": 2048,
            "temperature": 0.7,
        }

        resp = http_requests.post(url, headers=headers, json=payload, timeout=60)
        resp.raise_for_status()
        text = resp.json()["choices"][0]["message"]["content"]

        return self._parse_json_response(text, schema)

    def _parse_json_response(self, text: str, schema: ContentSchema) -> Dict[str, List[str]]:
        """Extract JSON from LLM response."""
        # Try to find JSON in code blocks
        json_match = re.search(r'```(?:json)?\s*\n?(.*?)\n?```', text, re.DOTALL)
        if json_match:
            text = json_match.group(1)

        try:
            result = json.loads(text.strip())
        except json.JSONDecodeError:
            # Fallback: use schema defaults
            result = {}
            for elem in schema.elements:
                if elem["type"] == "text":
                    result[elem["name"]] = elem.get("data", [""])
                else:
                    result[elem["name"]] = []

        # Ensure all elements present
        for elem in schema.elements:
            name = elem["name"]
            if name not in result:
                if elem["type"] == "text":
                    result[name] = elem.get("data", [""])
                else:
                    result[name] = []
            # Normalize to list
            if isinstance(result[name], str):
                result[name] = [result[name]]

        return result


# ═══════════════════════════════════════════════════════════════
#  XML-Level Slide Cloner
# ═══════════════════════════════════════════════════════════════

class SlideCloner:
    """
    Clone template slides at XML level and replace text.
    
    This is the core of the clone-based approach:
    - Deep copies slide XML (preserving all decorations, groups, images, formatting)
    - Replaces text in specific shapes by matching shape names/indices
    - Handles paragraph clone/delete for variable-length content
    - Preserves run formatting (font, color, size, bold, italic)
    """

    @staticmethod
    def clone_slide(prs: PptxPresentation, slide_idx: int) -> Any:
        """
        Clone a slide from the presentation by deep copying its XML.
        
        Args:
            prs: Source presentation
            slide_idx: 0-based index of the slide to clone
            
        Returns:
            New slide object added to the presentation
        """
        template_slide = prs.slides[slide_idx]
        layout = template_slide.slide_layout

        # Add a new slide with the same layout
        new_slide = prs.slides.add_slide(layout)

        # Remove all default shapes from the new slide
        for shape in list(new_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)

        # Deep copy all shape elements from template to new slide
        for shape in template_slide.shapes:
            el = deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        # Copy slide background if set
        if template_slide.background.fill.type is not None:
            bg_xml = deepcopy(template_slide.background._element)
            new_bg = new_slide.background._element
            new_bg.getparent().replace(new_bg, bg_xml)

        # Copy notes if any
        if template_slide.has_notes_slide:
            pass  # Notes are less critical

        return new_slide

    @staticmethod
    def replace_text_in_shape(shape, new_text: str, paragraph_idx: int = 0):
        """
        Replace text in a specific paragraph of a shape while preserving formatting.
        
        The key insight from PPTAgent: modify run text, don't recreate paragraphs.
        This preserves font, color, size, bold, italic, etc.
        
        Args:
            shape: python-pptx shape with text_frame
            new_text: New text to insert
            paragraph_idx: Which paragraph to replace (0-based)
        """
        if not shape.has_text_frame:
            return

        tf = shape.text_frame
        if paragraph_idx >= len(tf.paragraphs):
            return

        para = tf.paragraphs[paragraph_idx]

        if len(para.runs) == 0:
            # No runs — add text directly
            para.text = new_text
            return

        # Replace text in first run, clear the rest
        # This preserves the first run's formatting for all text
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""

    @staticmethod
    def replace_all_text_in_shape(shape, new_texts: List[str]):
        """
        Replace all paragraph text in a shape.
        
        Matches paragraphs 1:1. If new_texts has fewer items, keeps remaining
        paragraphs as-is. If more items, only fills up to existing paragraphs.
        
        Args:
            shape: python-pptx shape with text_frame
            new_texts: List of texts, one per paragraph
        """
        if not shape.has_text_frame:
            return

        tf = shape.text_frame
        for i, text in enumerate(new_texts):
            if i >= len(tf.paragraphs):
                break
            SlideCloner.replace_text_in_shape(shape, text, paragraph_idx=i)

    @staticmethod
    def clone_paragraph(shape, paragraph_idx: int):
        """
        Clone a paragraph within a shape (for adding more bullet points).
        
        Duplicates the XML of the specified paragraph and appends it.
        The cloned paragraph inherits all formatting.
        
        Args:
            shape: python-pptx shape with text_frame
            paragraph_idx: Index of the paragraph to clone
        """
        if not shape.has_text_frame:
            return

        tf = shape.text_frame
        if paragraph_idx >= len(tf.paragraphs):
            return

        source_para = tf.paragraphs[paragraph_idx]
        # Deep copy the paragraph XML element
        new_para_xml = deepcopy(source_para._element)
        # Insert after the last paragraph
        tf._txBody.append(new_para_xml)

    @staticmethod
    def delete_paragraph(shape, paragraph_idx: int):
        """
        Delete a paragraph from a shape.
        
        Keeps at least one paragraph (text frames require it).
        
        Args:
            shape: python-pptx shape with text_frame
            paragraph_idx: Index of the paragraph to delete
        """
        if not shape.has_text_frame:
            return

        tf = shape.text_frame
        if len(tf.paragraphs) <= 1:
            # Can't delete the last paragraph, just clear it
            tf.paragraphs[0].clear()
            return

        if paragraph_idx >= len(tf.paragraphs):
            return

        para = tf.paragraphs[paragraph_idx]
        para._element.getparent().remove(para._element)

    @staticmethod
    def adjust_paragraphs(shape, target_count: int, template_para_idx: int = 0):
        """
        Adjust paragraph count in a shape to match target.
        
        Clones or deletes paragraphs as needed.
        Uses the template_para_idx as the source for cloning.
        
        Args:
            shape: python-pptx shape with text_frame
            target_count: Desired number of paragraphs
            template_para_idx: Which paragraph to use as clone source
        """
        if not shape.has_text_frame:
            return

        tf = shape.text_frame
        current_count = len(tf.paragraphs)

        if target_count > current_count:
            # Need to clone paragraphs
            clone_idx = min(template_para_idx, current_count - 1)
            for _ in range(target_count - current_count):
                SlideCloner.clone_paragraph(shape, clone_idx)
        elif target_count < current_count:
            # Need to delete paragraphs (from the end)
            for _ in range(current_count - target_count):
                if len(tf.paragraphs) > 1:
                    SlideCloner.delete_paragraph(shape, len(tf.paragraphs) - 1)


# ═══════════════════════════════════════════════════════════════
#  Shape Matcher: Maps layout elements to actual shapes
# ═══════════════════════════════════════════════════════════════

class ShapeMatcher:
    """Match InducedLayout elements to actual slide shapes for editing."""

    @staticmethod
    def match_text_elements(slide, text_elements: List[TextElementInfo]) -> List[Tuple[Any, TextElementInfo]]:
        """
        Match text elements from induction to actual shapes in the cloned slide.
        
        Matching strategy (in order):
        1. By placeholder index (most reliable)
        2. By shape name (if unique)
        3. By position proximity
        
        Returns:
            List of (shape, TextElementInfo) pairs
        """
        matches = []
        used_shapes = set()

        # Build shape lookup
        text_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
                text_shapes.append(shape)

        # Also check group shapes recursively
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, 'shapes'):
                for child in shape.shapes:
                    if hasattr(child, 'text_frame') and child.has_text_frame and child.text_frame.text.strip():
                        text_shapes.append(child)

        # Pass 1: Match by placeholder idx
        for te in text_elements:
            if te.is_placeholder and te.placeholder_idx is not None:
                for shape in text_shapes:
                    if id(shape) in used_shapes:
                        continue
                    if (hasattr(shape, 'is_placeholder') and shape.is_placeholder and
                            hasattr(shape, 'placeholder_format') and
                            shape.placeholder_format.idx == te.placeholder_idx):
                        matches.append((shape, te))
                        used_shapes.add(id(shape))
                        break

        # Pass 2: Match by shape name
        matched_te_ids = {id(te) for _, te in matches}
        for te in text_elements:
            if id(te) in matched_te_ids:
                continue
            for shape in text_shapes:
                if id(shape) in used_shapes:
                    continue
                if shape.name == te.shape_name:
                    matches.append((shape, te))
                    used_shapes.add(id(shape))
                    matched_te_ids.add(id(te))
                    break

        # Pass 3: Match by position proximity
        for te in text_elements:
            if id(te) in matched_te_ids:
                continue
            best_shape = None
            best_dist = float('inf')
            for shape in text_shapes:
                if id(shape) in used_shapes:
                    continue
                # Distance in EMU space
                dx = abs(shape.left - te.left)
                dy = abs(shape.top - te.top)
                dist = dx + dy
                if dist < best_dist:
                    best_dist = dist
                    best_shape = shape
            if best_shape and best_dist < 914400 * 2:  # Within ~2 inches
                matches.append((best_shape, te))
                used_shapes.add(id(best_shape))

        return matches


# ═══════════════════════════════════════════════════════════════
#  Clone Generator: Main orchestrator
# ═══════════════════════════════════════════════════════════════

class CloneGenerator:
    """
    Generate presentation by cloning template slides and replacing text.
    
    Workflow:
    1. TemplateInducter.induct() → TemplateSpec
    2. For each output slide:
       a. Select layout from TemplateSpec
       b. ContentFiller fills schema with LLM
       c. SlideCloner clones the template slide
       d. ShapeMatcher maps elements to shapes
       e. Text replacement + paragraph adjustment
    3. Save output PPTX
    
    Usage:
        gen = CloneGenerator()
        spec = gen.induct_template("template.pptx")
        output_path = gen.generate(
            spec=spec,
            outline=[
                {"purpose": "Cover", "topic": "My Presentation"},
                {"purpose": "Content", "content": "Key findings about..."},
                ...
            ],
            output_path="output.pptx"
        )
    """

    def __init__(self):
        self.inducter = TemplateInducter()
        self.filler = ContentFiller()

    def induct_template(self, pptx_path: str) -> TemplateSpec:
        """Run induction on a template file."""
        return self.inducter.induct(pptx_path)

    def generate(
        self,
        spec: TemplateSpec,
        outline: List[Dict[str, Any]],
        output_path: str,
        source_content: str = "",
    ) -> str:
        """
        Generate a full presentation using clone-based approach.
        
        Args:
            spec: TemplateSpec from induction
            outline: List of slide descriptions, each with:
                - purpose: str (e.g., "Cover", "Content", "Section", "Ending")
                - topic: str (slide title/subject)
                - content: str (source content for this slide, optional)
            output_path: Path to save the output PPTX
            source_content: Global source content (used if individual slides don't have content)
            
        Returns:
            Absolute path to the generated file
        """
        # Open the template
        prs = PptxPresentation(spec.source_path)

        # We'll build a new presentation by cloning slides
        # Strategy: work with the existing prs, clone slides, then remove originals

        # Count original slides
        original_count = len(prs.slides)

        # Generate each slide
        for slide_info in outline:
            purpose = slide_info.get("purpose", "Content")
            topic = slide_info.get("topic", "")
            content = slide_info.get("content", source_content)

            # Select layout
            layout = self._select_layout(spec, purpose)
            if layout is None:
                # Fallback: use first content layout
                layout = spec.content_layouts[0] if spec.content_layouts else spec.layouts[0]

            # Generate content for schema
            filled_content = self._fill_content(layout, topic, content)

            # Clone the template slide
            new_slide = SlideCloner.clone_slide(prs, layout.template_slide_idx)

            # Replace text
            self._apply_content(new_slide, layout, filled_content)

        # Remove original template slides
        self._remove_original_slides(prs, original_count)

        # Save
        out_path = Path(output_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))

        return str(out_path.absolute())

    def generate_single_slide(
        self,
        spec: TemplateSpec,
        prs: PptxPresentation,
        purpose: str,
        topic: str,
        content: str,
    ) -> Any:
        """
        Generate a single slide and add it to the presentation.
        
        Useful for streaming/incremental generation.
        
        Args:
            spec: TemplateSpec
            prs: Target presentation to add slide to
            purpose: Slide purpose
            topic: Slide topic
            content: Source content
            
        Returns:
            The new slide object
        """
        layout = self._select_layout(spec, purpose)
        if layout is None:
            layout = spec.content_layouts[0] if spec.content_layouts else spec.layouts[0]

        filled_content = self._fill_content(layout, topic, content)
        new_slide = SlideCloner.clone_slide(prs, layout.template_slide_idx)
        self._apply_content(new_slide, layout, filled_content)
        return new_slide

    # ─────────────────────────────────────────────────
    #  Internal Methods
    # ─────────────────────────────────────────────────

    def _select_layout(self, spec: TemplateSpec, purpose: str) -> Optional[InducedLayout]:
        """Select the best layout for a given slide purpose."""
        purpose_lower = purpose.lower()

        # Map purposes to categories
        category_map = {
            "cover": SlideCategory.COVER,
            "title": SlideCategory.COVER,
            "opening": SlideCategory.COVER,
            "toc": SlideCategory.TOC,
            "table of contents": SlideCategory.TOC,
            "agenda": SlideCategory.TOC,
            "section": SlideCategory.SECTION,
            "section outline": SlideCategory.SECTION,
            "divider": SlideCategory.SECTION,
            "ending": SlideCategory.ENDING,
            "thank": SlideCategory.ENDING,
            "summary": SlideCategory.ENDING,
            "conclusion": SlideCategory.ENDING,
        }

        for keyword, category in category_map.items():
            if keyword in purpose_lower:
                layout = spec.get_layout_by_category(category)
                if layout:
                    return layout

        # Default: content layout
        has_image = "image" in purpose_lower or "图" in purpose_lower
        return spec.get_content_layout(has_image=has_image)

    def _fill_content(self, layout: InducedLayout,
                      topic: str, content: str) -> Dict[str, List[str]]:
        """Fill content schema using LLM."""
        if not layout.schema.elements:
            return {}

        try:
            return self.filler.fill(layout.schema, topic, content)
        except Exception as e:
            print(f"  ⚠ ContentFiller failed: {e}, using defaults")
            # Fallback: use schema defaults
            result = {}
            for elem in layout.schema.elements:
                if elem["type"] == "text":
                    result[elem["name"]] = elem.get("data", [topic])
                else:
                    result[elem["name"]] = []
            return result

    def _apply_content(self, slide, layout: InducedLayout,
                       filled_content: Dict[str, List[str]]):
        """Apply filled content to a cloned slide."""
        # Match shapes to layout elements
        matches = ShapeMatcher.match_text_elements(slide, layout.text_elements)

        # Create a mapping from TextElementInfo role/name to filled content
        role_content_map = self._build_role_content_map(layout, filled_content)

        for shape, te_info in matches:
            content_key = self._find_content_key(te_info, layout, filled_content)
            if content_key and content_key in filled_content:
                texts = filled_content[content_key]
                if not texts:
                    continue

                if te_info.paragraph_count > 1 and len(texts) > 1:
                    # Multi-paragraph element (e.g., bullet list)
                    SlideCloner.adjust_paragraphs(shape, len(texts))
                    SlideCloner.replace_all_text_in_shape(shape, texts)
                else:
                    # Single paragraph element
                    SlideCloner.replace_text_in_shape(shape, texts[0])

    def _build_role_content_map(self, layout: InducedLayout,
                                filled_content: Dict[str, List[str]]) -> Dict[str, str]:
        """Map text element roles to content keys."""
        role_map = {}
        for elem in layout.schema.elements:
            if elem["type"] == "text":
                role_map[elem["name"]] = elem["name"]
        return role_map

    def _find_content_key(self, te_info: TextElementInfo,
                          layout: InducedLayout,
                          filled_content: Dict[str, List[str]]) -> Optional[str]:
        """Find the content key for a text element."""
        # Direct name match
        for elem in layout.schema.elements:
            if elem["type"] != "text":
                continue
            name = elem["name"]
            if name in filled_content:
                # Match by role
                if te_info.role == name or te_info.role in name:
                    return name
                # Match by element naming convention
                if (te_info.role == "title" and "title" in name) or \
                   (te_info.role == "subtitle" and "subtitle" in name) or \
                   (te_info.role == "body" and "body" in name) or \
                   (te_info.role == "label" and "label" in name):
                    return name

        # Fallback: first content key matching role prefix
        for key in filled_content:
            if te_info.role in key:
                return key

        return None

    def _remove_original_slides(self, prs: PptxPresentation, count: int):
        """Remove the first `count` slides (original template slides)."""
        for _ in range(count):
            if len(prs.slides._sldIdLst) > count:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]


# ═══════════════════════════════════════════════════════════════
#  Convenience: Full Pipeline Function
# ═══════════════════════════════════════════════════════════════

def clone_generate_pptx(
    template_path: str,
    outline: List[Dict[str, Any]],
    output_path: str,
    source_content: str = "",
) -> str:
    """
    One-shot: induct template + generate presentation.
    
    Args:
        template_path: Path to the PPTX template
        outline: List of slide descriptions
        output_path: Where to save the result
        source_content: Global source content
        
    Returns:
        Path to the generated PPTX
    """
    gen = CloneGenerator()
    spec = gen.induct_template(template_path)
    return gen.generate(spec, outline, output_path, source_content)
