"""
Clone-based Template Inducter for LivePPT.

Inspired by PPTAgent's SlideInducter but adapted for text-only LLM:
- Uses python-pptx shape analysis (no vision model) for layout clustering
- Heuristic-based functional slide classification (no multimodal)
- Deep shape extraction including Groups (recursive)
- Generates content schemas for LLM-driven content filling

Two-phase approach:
  Phase 1 (Induction): PPTX → TemplateSpec (layouts, schemas, functional slides)
  Phase 2 (Generation): TemplateSpec + content → cloned PPTX (see clone_generator.py)
"""

import re
import hashlib
from copy import deepcopy
from collections import defaultdict
from dataclasses import dataclass, field
from typing import List, Dict, Optional, Tuple, Any
from pathlib import Path
from enum import Enum

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ═══════════════════════════════════════════════════════════════
#  Data Models
# ═══════════════════════════════════════════════════════════════

class SlideCategory(str, Enum):
    """Functional slide categories."""
    COVER = "cover"
    TOC = "toc"
    SECTION = "section"
    ENDING = "ending"
    CONTENT = "content"


@dataclass
class TextElementInfo:
    """Extracted text element from a shape."""
    shape_idx: int
    shape_name: str
    text: str
    font_size_pt: Optional[float] = None
    font_name: Optional[str] = None
    font_bold: Optional[bool] = None
    font_color: Optional[str] = None
    is_placeholder: bool = False
    placeholder_idx: Optional[int] = None
    paragraph_count: int = 1
    # Position (EMU values stored for building)
    left: int = 0
    top: int = 0
    width: int = 0
    height: int = 0
    # Role classification
    role: str = "body"  # "title", "subtitle", "body", "label", "number"

    @property
    def char_count(self) -> int:
        return len(self.text)


@dataclass
class ImageElementInfo:
    """Extracted image element from a shape."""
    shape_idx: int
    shape_name: str
    left: int = 0
    top: int = 0
    width: int = 0
    height: int = 0
    image_sha1: Optional[str] = None
    image_ext: Optional[str] = None


@dataclass
class ContentSchema:
    """Schema describing what content a layout expects.
    
    Each element has a name, type (text/image), and sample data.
    The LLM uses this to fill content into cloned slides.
    """
    elements: List[Dict[str, Any]] = field(default_factory=list)

    def add_text_element(self, name: str, sample_data: List[str],
                         suggested_chars: Optional[int] = None,
                         variable_length: Optional[Tuple[int, int]] = None):
        elem = {
            "name": name,
            "type": "text",
            "data": sample_data,
            "suggested_characters": suggested_chars or max((len(s) for s in sample_data), default=20),
        }
        if variable_length:
            elem["variable_length"] = variable_length
        self.elements.append(elem)

    def add_image_element(self, name: str, count: int = 1):
        self.elements.append({
            "name": name,
            "type": "image",
            "count": count,
        })

    def to_prompt(self) -> str:
        """Generate a schema description for LLM prompts."""
        lines = []
        for elem in self.elements:
            line = f"Element: {elem['name']}\n  type: {elem['type']}"
            if elem["type"] == "text":
                line += f"\n  suggested_characters: {elem.get('suggested_characters', 50)}"
                line += f"\n  default_quantity: {len(elem.get('data', []))}"
                if elem.get("variable_length"):
                    vl = elem["variable_length"]
                    line += f"\n  variable_length: [{vl[0]}, {vl[1]}]"
            elif elem["type"] == "image":
                line += f"\n  count: {elem.get('count', 1)}"
            lines.append(line)
        return "\n".join(lines)

    def to_dict(self) -> List[Dict[str, Any]]:
        return self.elements


@dataclass
class InducedLayout:
    """A layout pattern induced from template slides.
    
    Attributes:
        name: Layout name (e.g., "content:text", "cover", "section")
        category: Functional category
        template_slide_idx: 0-based index of the representative slide in the template
        slide_indices: All slide indices that match this layout
        schema: Content schema describing what can be filled
        text_elements: Extracted text element info
        image_elements: Extracted image element info
        shape_signature: Hash-based signature for clustering
    """
    name: str
    category: SlideCategory
    template_slide_idx: int
    slide_indices: List[int] = field(default_factory=list)
    schema: ContentSchema = field(default_factory=ContentSchema)
    text_elements: List[TextElementInfo] = field(default_factory=list)
    image_elements: List[ImageElementInfo] = field(default_factory=list)
    shape_signature: str = ""

    @property
    def has_images(self) -> bool:
        return len(self.image_elements) > 0

    def to_dict(self) -> Dict[str, Any]:
        return {
            "name": self.name,
            "category": self.category.value,
            "template_slide_idx": self.template_slide_idx,
            "slide_indices": self.slide_indices,
            "schema": self.schema.to_dict(),
            "num_text_elements": len(self.text_elements),
            "num_image_elements": len(self.image_elements),
        }


@dataclass
class TemplateSpec:
    """Complete template specification from induction.
    
    This is the output of TemplateInducter.induct() and the input to CloneGenerator.
    """
    source_path: str
    slide_width: float  # in points
    slide_height: float  # in points
    num_slides: int
    layouts: List[InducedLayout] = field(default_factory=list)
    functional_layouts: Dict[str, InducedLayout] = field(default_factory=dict)
    content_layouts: List[InducedLayout] = field(default_factory=list)

    def get_layout_by_category(self, category: SlideCategory) -> Optional[InducedLayout]:
        return self.functional_layouts.get(category.value)

    def get_content_layout(self, has_image: bool = False) -> Optional[InducedLayout]:
        """Get a content layout, preferring image layouts if has_image."""
        if has_image:
            img_layouts = [l for l in self.content_layouts if l.has_images]
            if img_layouts:
                return img_layouts[0]
        text_layouts = [l for l in self.content_layouts if not l.has_images]
        return text_layouts[0] if text_layouts else (self.content_layouts[0] if self.content_layouts else None)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "source_path": self.source_path,
            "slide_width_pt": self.slide_width,
            "slide_height_pt": self.slide_height,
            "num_slides": self.num_slides,
            "layouts": [l.to_dict() for l in self.layouts],
            "functional_layout_names": list(self.functional_layouts.keys()),
            "content_layout_count": len(self.content_layouts),
        }


# ═══════════════════════════════════════════════════════════════
#  Template Inducter
# ═══════════════════════════════════════════════════════════════

class TemplateInducter:
    """
    Analyze a PPTX template and induce layout patterns.
    
    Adapted from PPTAgent's SlideInducter for text-only LLM:
    - category_split: Heuristic rules (font size, position, keywords) instead of vision model
    - layout_split: Shape signature clustering instead of image embeddings
    - content_induct: Extracts text/image elements and builds content schemas
    
    Usage:
        inducter = TemplateInducter()
        spec = inducter.induct("template.pptx")
        # spec.layouts, spec.functional_layouts, spec.content_layouts
    """

    # Heuristic thresholds
    TITLE_FONT_SIZE_MIN = 24  # pt
    LARGE_TITLE_FONT_SIZE = 36  # pt
    COVER_KEYWORDS = {"welcome", "introduction", "presented by", "欢迎", "汇报",
                      "报告", "演示", "presentation", "title"}
    TOC_KEYWORDS = {"contents", "table of contents", "agenda", "outline", "目录",
                    "大纲", "议程", "content"}
    SECTION_KEYWORDS = {"section", "part", "chapter", "章节", "部分", "篇"}
    ENDING_KEYWORDS = {"thank", "thanks", "q&a", "questions", "结束", "谢谢",
                       "感谢", "end", "summary", "总结"}

    def induct(self, pptx_path: str) -> TemplateSpec:
        """
        Analyze template and produce a TemplateSpec.
        
        Args:
            pptx_path: Path to the PPTX template file
            
        Returns:
            TemplateSpec with all induced layouts
        """
        path = Path(pptx_path)
        if not path.exists():
            raise FileNotFoundError(f"Template not found: {pptx_path}")

        prs = Presentation(str(path))
        slide_width_pt = prs.slide_width.pt if hasattr(prs.slide_width, 'pt') else prs.slide_width / 914400 * 72
        slide_height_pt = prs.slide_height.pt if hasattr(prs.slide_height, 'pt') else prs.slide_height / 914400 * 72
        slide_area = slide_width_pt * slide_height_pt

        # Step 1: Extract all slides
        slide_infos = []
        for idx, slide in enumerate(prs.slides):
            info = self._extract_slide_info(slide, idx, slide_width_pt, slide_height_pt)
            slide_infos.append(info)

        # Step 2: Classify functional slides (heuristic)
        categories = self._classify_slides(slide_infos)

        # Step 3: Cluster content slides by shape signature
        content_indices = [i for i, cat in enumerate(categories) if cat == SlideCategory.CONTENT]
        clusters = self._cluster_by_signature(slide_infos, content_indices)

        # Step 4: Build layouts
        layouts = []
        functional_layouts = {}
        content_layouts = []

        # Functional layouts
        for idx, cat in enumerate(categories):
            if cat == SlideCategory.CONTENT:
                continue
            layout = self._build_layout(
                name=cat.value,
                category=cat,
                slide_infos=slide_infos,
                representative_idx=idx,
                member_indices=[idx],
            )
            layouts.append(layout)
            functional_layouts[cat.value] = layout

        # Content layouts (from clusters)
        for cluster_id, member_indices in enumerate(clusters):
            representative = member_indices[0]
            info = slide_infos[representative]
            has_images = info["num_images"] > 0
            name = f"content:{'image' if has_images else 'text'}:{cluster_id}"

            layout = self._build_layout(
                name=name,
                category=SlideCategory.CONTENT,
                slide_infos=slide_infos,
                representative_idx=representative,
                member_indices=member_indices,
            )
            layouts.append(layout)
            content_layouts.append(layout)

        spec = TemplateSpec(
            source_path=str(path.absolute()),
            slide_width=slide_width_pt,
            slide_height=slide_height_pt,
            num_slides=len(prs.slides),
            layouts=layouts,
            functional_layouts=functional_layouts,
            content_layouts=content_layouts,
        )

        return spec

    # ─────────────────────────────────────────────────
    #  Step 1: Extract Slide Info
    # ─────────────────────────────────────────────────

    def _extract_slide_info(self, slide, slide_idx: int,
                            slide_w_pt: float, slide_h_pt: float) -> Dict[str, Any]:
        """Extract structured info from a slide for classification and clustering."""
        text_elements = []
        image_elements = []
        shape_types = []
        all_text = ""

        shape_idx = 0
        for shape in slide.shapes:
            self._extract_shape_recursive(
                shape, slide_idx, shape_idx,
                text_elements, image_elements, shape_types,
                slide_w_pt, slide_h_pt
            )
            shape_idx += 1

        all_text = " ".join(te.text for te in text_elements).lower()

        # Classify text roles
        self._classify_text_roles(text_elements, slide_w_pt, slide_h_pt)

        # Build shape signature for clustering
        signature = self._compute_shape_signature(text_elements, image_elements, shape_types)

        return {
            "slide_idx": slide_idx,
            "text_elements": text_elements,
            "image_elements": image_elements,
            "shape_types": shape_types,
            "all_text": all_text,
            "num_text": len(text_elements),
            "num_images": len(image_elements),
            "signature": signature,
            "title_text": self._find_title_text(text_elements),
        }

    def _extract_shape_recursive(self, shape, slide_idx: int, shape_idx: int,
                                 text_elements: List[TextElementInfo],
                                 image_elements: List[ImageElementInfo],
                                 shape_types: List[str],
                                 slide_w_pt: float, slide_h_pt: float,
                                 depth: int = 0):
        """Recursively extract shapes, handling groups."""
        if depth > 5:  # Prevent infinite recursion
            return

        shape_type = shape.shape_type if hasattr(shape, 'shape_type') else None

        # Group shape: recurse into children
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            shape_types.append("group")
            if hasattr(shape, 'shapes'):
                for child_idx, child in enumerate(shape.shapes):
                    self._extract_shape_recursive(
                        child, slide_idx, shape_idx * 100 + child_idx,
                        text_elements, image_elements, shape_types,
                        slide_w_pt, slide_h_pt, depth + 1
                    )
            return

        # Picture
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_types.append("picture")
            img_info = ImageElementInfo(
                shape_idx=shape_idx,
                shape_name=shape.name,
                left=shape.left,
                top=shape.top,
                width=shape.width,
                height=shape.height,
            )
            if hasattr(shape, 'image'):
                img_info.image_sha1 = shape.image.sha1
                img_info.image_ext = shape.image.ext
            image_elements.append(img_info)
            return

        # Text-bearing shape
        if hasattr(shape, 'text_frame') and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Get dominant font info
                font_size, font_name, font_bold, font_color = self._get_dominant_font(shape)

                is_ph = shape.is_placeholder if hasattr(shape, 'is_placeholder') else False
                ph_idx = shape.placeholder_format.idx if (is_ph and hasattr(shape, 'placeholder_format')) else None

                te = TextElementInfo(
                    shape_idx=shape_idx,
                    shape_name=shape.name,
                    text=text,
                    font_size_pt=font_size,
                    font_name=font_name,
                    font_bold=font_bold,
                    font_color=font_color,
                    is_placeholder=is_ph,
                    placeholder_idx=ph_idx,
                    paragraph_count=len(shape.text_frame.paragraphs),
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                )
                text_elements.append(te)

                st = str(shape_type).split("(")[0].lower() if shape_type else "textbox"
                shape_types.append(st)
            else:
                shape_types.append("empty_text")
        elif shape_type == MSO_SHAPE_TYPE.TABLE:
            shape_types.append("table")
        elif shape_type == MSO_SHAPE_TYPE.CHART:
            shape_types.append("chart")
        elif shape_type == MSO_SHAPE_TYPE.FREEFORM:
            shape_types.append("freeform")
        elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_types.append("autoshape")
        else:
            shape_types.append("other")

    def _get_dominant_font(self, shape) -> Tuple[Optional[float], Optional[str], Optional[bool], Optional[str]]:
        """Get the most common/largest font properties from a shape's text frame."""
        if not hasattr(shape, 'text_frame'):
            return None, None, None, None

        max_size = None
        font_name = None
        font_bold = None
        font_color = None

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    size_pt = run.font.size.pt if hasattr(run.font.size, 'pt') else None
                    if size_pt and (max_size is None or size_pt > max_size):
                        max_size = size_pt
                        font_name = run.font.name
                        font_bold = run.font.bold
                        try:
                            if run.font.color and run.font.color.type is not None:
                                font_color = str(run.font.color.rgb)
                        except (AttributeError, TypeError):
                            pass  # Scheme color, theme color, etc.
        return max_size, font_name, font_bold, font_color

    def _classify_text_roles(self, text_elements: List[TextElementInfo],
                             slide_w_pt: float, slide_h_pt: float):
        """Classify each text element's role based on font size and position."""
        if not text_elements:
            return

        # Sort by font size (descending), then by vertical position
        for te in text_elements:
            size = te.font_size_pt or 12
            top_ratio = te.top / (slide_h_pt * 12700) if slide_h_pt else 0.5  # EMU to ratio

            if size >= self.LARGE_TITLE_FONT_SIZE:
                te.role = "title"
            elif size >= self.TITLE_FONT_SIZE_MIN:
                te.role = "subtitle" if any(
                    t.role == "title" for t in text_elements if t is not te
                ) else "title"
            elif te.paragraph_count >= 3:
                te.role = "body"
            elif te.char_count <= 10 and size and size <= 14:
                te.role = "label"
            else:
                te.role = "body"

    def _find_title_text(self, text_elements: List[TextElementInfo]) -> str:
        """Find the title text from extracted elements."""
        titles = [te for te in text_elements if te.role == "title"]
        if titles:
            return titles[0].text
        # Fallback: largest font
        if text_elements:
            by_size = sorted(text_elements, key=lambda t: t.font_size_pt or 0, reverse=True)
            return by_size[0].text
        return ""

    # ─────────────────────────────────────────────────
    #  Step 2: Classify Functional Slides (Heuristic)
    # ─────────────────────────────────────────────────

    def _classify_slides(self, slide_infos: List[Dict]) -> List[SlideCategory]:
        """Classify each slide as functional or content using heuristics."""
        n = len(slide_infos)
        categories = [SlideCategory.CONTENT] * n

        if n == 0:
            return categories

        # Rule 1: First slide is likely cover
        if n >= 1:
            info = slide_infos[0]
            if self._is_likely_cover(info, is_first=True):
                categories[0] = SlideCategory.COVER

        # Rule 2: Last slide is likely ending
        if n >= 2:
            info = slide_infos[-1]
            if self._is_likely_ending(info, is_last=True):
                categories[-1] = SlideCategory.ENDING

        # Rule 3: Check all slides for TOC and section dividers
        for i, info in enumerate(slide_infos):
            if categories[i] != SlideCategory.CONTENT:
                continue

            if self._is_likely_toc(info):
                categories[i] = SlideCategory.TOC
            elif self._is_likely_section(info, i, n):
                categories[i] = SlideCategory.SECTION

        return categories

    def _is_likely_cover(self, info: Dict, is_first: bool = False) -> bool:
        """Check if slide is likely a cover/title page."""
        text = info["all_text"]
        # Keyword match
        if any(kw in text for kw in self.COVER_KEYWORDS):
            return True
        # First slide with large title and few elements
        if is_first:
            titles = [te for te in info["text_elements"] if te.role == "title"]
            if titles and info["num_text"] <= 4:
                return True
        return False

    def _is_likely_ending(self, info: Dict, is_last: bool = False) -> bool:
        """Check if slide is likely an ending/thank you page."""
        text = info["all_text"]
        if any(kw in text for kw in self.ENDING_KEYWORDS):
            return True
        if is_last and info["num_text"] <= 2:
            return True
        return False

    def _is_likely_toc(self, info: Dict) -> bool:
        """Check if slide is likely a table of contents."""
        text = info["all_text"]
        return any(kw in text for kw in self.TOC_KEYWORDS)

    def _is_likely_section(self, info: Dict, idx: int, total: int) -> bool:
        """Check if slide is likely a section divider."""
        text = info["all_text"]
        # Keyword match
        if any(kw in text for kw in self.SECTION_KEYWORDS):
            return True
        # Heuristic: very few text elements, one large title, no body text
        if info["num_text"] <= 2 and info["num_images"] == 0:
            titles = [te for te in info["text_elements"] if te.role == "title"]
            bodies = [te for te in info["text_elements"] if te.role == "body"]
            if len(titles) >= 1 and len(bodies) == 0:
                # Could be a section divider
                return True
        return False

    # ─────────────────────────────────────────────────
    #  Step 3: Shape Signature Clustering
    # ─────────────────────────────────────────────────

    def _compute_shape_signature(self, text_elements: List[TextElementInfo],
                                 image_elements: List[ImageElementInfo],
                                 shape_types: List[str]) -> str:
        """Compute a signature for layout clustering.
        
        Encodes: shape type counts, text role distribution, position grid.
        """
        parts = []

        # Shape type counts
        type_counts = defaultdict(int)
        for st in shape_types:
            type_counts[st] += 1
        type_str = ",".join(f"{k}:{v}" for k, v in sorted(type_counts.items()))
        parts.append(type_str)

        # Text role distribution
        role_counts = defaultdict(int)
        for te in text_elements:
            role_counts[te.role] += 1
        role_str = ",".join(f"{k}:{v}" for k, v in sorted(role_counts.items()))
        parts.append(role_str)

        # Number of images
        parts.append(f"img:{len(image_elements)}")

        sig_input = "|".join(parts)
        return hashlib.md5(sig_input.encode()).hexdigest()[:8]

    def _cluster_by_signature(self, slide_infos: List[Dict],
                              content_indices: List[int]) -> List[List[int]]:
        """Cluster content slides by their shape signature."""
        sig_groups = defaultdict(list)
        for idx in content_indices:
            sig = slide_infos[idx]["signature"]
            sig_groups[sig].append(idx)

        # Return clusters sorted by first occurrence
        clusters = sorted(sig_groups.values(), key=lambda x: x[0])
        return clusters

    # ─────────────────────────────────────────────────
    #  Step 4: Build InducedLayout
    # ─────────────────────────────────────────────────

    def _build_layout(self, name: str, category: SlideCategory,
                      slide_infos: List[Dict],
                      representative_idx: int,
                      member_indices: List[int]) -> InducedLayout:
        """Build an InducedLayout from the representative slide."""
        info = slide_infos[representative_idx]
        text_elements = info["text_elements"]
        image_elements = info["image_elements"]

        # Build content schema
        schema = ContentSchema()

        # Add text elements grouped by role
        role_groups = defaultdict(list)
        for te in text_elements:
            role_groups[te.role].append(te)

        for role in ["title", "subtitle", "body", "label", "number"]:
            elements = role_groups.get(role, [])
            if not elements:
                continue

            if role in ("title", "subtitle"):
                # Single text element
                for i, te in enumerate(elements):
                    elem_name = role if i == 0 else f"{role}_{i+1}"
                    schema.add_text_element(
                        name=elem_name,
                        sample_data=[te.text],
                        suggested_chars=te.char_count,
                    )
            elif role == "body":
                # Body elements may have variable length (bullet points)
                for i, te in enumerate(elements):
                    elem_name = f"body_{i}" if len(elements) > 1 else "body"
                    paragraphs = te.text.split("\n")
                    paragraphs = [p.strip() for p in paragraphs if p.strip()]
                    if te.paragraph_count > 2:
                        # Variable-length list element
                        schema.add_text_element(
                            name=elem_name,
                            sample_data=paragraphs,
                            suggested_chars=max((len(p) for p in paragraphs), default=30),
                            variable_length=(1, max(te.paragraph_count + 2, 8)),
                        )
                    else:
                        schema.add_text_element(
                            name=elem_name,
                            sample_data=paragraphs or [te.text],
                            suggested_chars=te.char_count,
                        )
            else:
                for i, te in enumerate(elements):
                    elem_name = f"{role}_{i}" if len(elements) > 1 else role
                    schema.add_text_element(
                        name=elem_name,
                        sample_data=[te.text],
                        suggested_chars=te.char_count,
                    )

        # Add image elements
        if image_elements:
            schema.add_image_element("images", count=len(image_elements))

        layout = InducedLayout(
            name=name,
            category=category,
            template_slide_idx=representative_idx,
            slide_indices=member_indices,
            schema=schema,
            text_elements=text_elements,
            image_elements=image_elements,
            shape_signature=info["signature"],
        )

        return layout
