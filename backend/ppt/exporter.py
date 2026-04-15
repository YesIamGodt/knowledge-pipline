"""
PPTX file exporter for LivePPT.

Converts slide JSON to python-pptx Presentation objects.
Two modes:
  - Freeform: 7 built-in slide types, absolute positioning (no template)
  - Template: Uses uploaded .pptx template layouts + placeholder filling

Freeform supports: title, bullets, comparison, metric, quote, timeline, flowchart.
"""

import base64
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional, Any

try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from backend.ppt.models import Slide, SlideType, TemplateStyle


# Default colors
DARK_BG = RGBColor(0x0F, 0x13, 0x18) if HAS_PPTX else None
LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF) if HAS_PPTX else None
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF) if HAS_PPTX else None
MUTED_TEXT = RGBColor(0x8B, 0x94, 0x9E) if HAS_PPTX else None
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
PX_TO_IN_X = SLIDE_W_IN / 800
PX_TO_IN_Y = SLIDE_H_IN / 600


class PPTXExporter:
    """Export Slide models to .pptx files."""

    def __init__(self, template_style: Optional[TemplateStyle] = None):
        if not HAS_PPTX:
            raise ImportError("python-pptx is required. Install: pip install python-pptx")
        self.style = template_style

    def export(
        self,
        slides: List[Dict[str, Any]],
        output_path: str,
        title: str = "Live PPT",
        template_path: Optional[str] = None,
    ) -> str:
        """
        Export slides to a .pptx file.

        Args:
            slides: List of slide dicts (as returned by LLM)
            output_path: Output file path
            title: Presentation title
            template_path: Optional .pptx template file to use as base

        Returns:
            Absolute path to the generated file
        """
        if template_path and Path(template_path).exists():
            return self._export_with_template(slides, output_path, template_path)

        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides:
            slide_type = slide_data.get("type", "bullets")
            self._add_slide(prs, slide_data, slide_type)

        out_path = Path(output_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        return str(out_path.absolute())

    def _add_slide(self, prs: PptxPresentation, data: Dict, slide_type: str):
        """Add one slide to the presentation."""
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)

        # Handle HTML-format slides (from LLM HTML mode)
        if data.get("html"):
            self._render_html_slide(slide, data)
            return

        theme = data.get("theme", {}) if isinstance(data.get("theme"), dict) else {}

        # Background color
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._parse_color(theme.get("bg"), DARK_BG)

        if isinstance(data.get("elements"), list) and data.get("elements"):
            self._render_elements_slide(slide, data)
            return

        method_map = {
            "title": self._render_title,
            "bullets": self._render_bullets,
            "comparison": self._render_comparison,
            "metric": self._render_metric,
            "quote": self._render_quote,
            "timeline": self._render_timeline,
            "flowchart": self._render_flowchart,
        }
        renderer = method_map.get(slide_type, self._render_bullets)
        renderer(slide, data)

        # Optional user-added image from live editing
        self._add_optional_image(slide, data)

        # Source annotation
        source = data.get("source", "")
        if source:
            self._add_source(slide, source)

    def _render_html_slide(self, slide, data: Dict):
        """Render an HTML-format slide by parsing inline styles and extracting text."""
        import re
        from html.parser import HTMLParser

        html = data["html"]

        # ── Extract background color from root div ──
        bg_color = DARK_BG
        # Match the first background or background-color in style attribute
        root_style_match = re.search(r'<div[^>]*style="([^"]*)"', html)
        if root_style_match:
            style_str = root_style_match.group(1)
            bg_match = re.search(r'background(?:-color)?\s*:\s*([^;]+)', style_str)
            if bg_match:
                bg_val = bg_match.group(1).strip()
                # Extract first hex color (works for solid and gradient)
                hex_colors = re.findall(r'#([0-9a-fA-F]{6})', bg_val)
                if hex_colors:
                    h = hex_colors[0]
                    try:
                        bg_color = RGBColor(
                            int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                        )
                    except (ValueError, IndexError):
                        pass

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # ── Parse HTML to extract text segments with styles ──
        class _TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.segments = []
                self.style_stack = [{}]

            def handle_starttag(self, tag, attrs):
                parent_style = dict(self.style_stack[-1])
                for attr, val in attrs:
                    if attr == 'style' and val:
                        for part in val.split(';'):
                            if ':' in part:
                                k, v = part.split(':', 1)
                                parent_style[k.strip().lower()] = v.strip()
                self.style_stack.append(parent_style)

            def handle_endtag(self, tag):
                if len(self.style_stack) > 1:
                    self.style_stack.pop()

            def handle_data(self, data_text):
                text = data_text.strip()
                if text:
                    style = dict(self.style_stack[-1]) if self.style_stack else {}
                    self.segments.append((text, style))

        parser = _TextExtractor()
        parser.feed(html)

        if not parser.segments:
            return

        # ── Lay out text segments on the slide ──
        y_offset = 0.4  # inches from top

        for text, style in parser.segments:
            if len(text.strip()) < 1:
                continue

            # Parse font size (px → pt ≈ px * 0.75)
            fs_px = 14
            fs_val = style.get('font-size', '')
            m = re.search(r'(\d+)', fs_val)
            if m:
                fs_px = int(m.group(1))
            pt_size = max(8, min(fs_px * 0.75, 44))

            # Parse text color
            text_color = LIGHT_TEXT
            color_val = style.get('color', '')
            hex_m = re.search(r'#([0-9a-fA-F]{6})', color_val)
            if hex_m:
                h = hex_m.group(1)
                try:
                    text_color = RGBColor(
                        int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                    )
                except (ValueError, IndexError):
                    pass
            else:
                rgb_m = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', color_val)
                if rgb_m:
                    try:
                        text_color = RGBColor(
                            int(rgb_m.group(1)),
                            int(rgb_m.group(2)),
                            int(rgb_m.group(3)),
                        )
                    except (ValueError, IndexError):
                        pass

            # Skip very low-opacity text
            opacity_val = style.get('opacity', '1')
            try:
                if float(opacity_val) < 0.15:
                    continue
            except ValueError:
                pass

            # Bold detection
            fw = style.get('font-weight', '')
            bold = fw in ('bold', '600', '700', '800', '900')

            # Height estimate
            h_estimate = max(0.4, pt_size / 72 * 1.8)

            # Don't overflow slide
            if y_offset + h_estimate > SLIDE_H_IN - 0.3:
                break

            # Alignment
            align = PP_ALIGN.LEFT
            ta = style.get('text-align', '')
            if ta == 'center':
                align = PP_ALIGN.CENTER
            elif ta == 'right':
                align = PP_ALIGN.RIGHT

            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_offset),
                Inches(SLIDE_W_IN - 1), Inches(h_estimate),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(int(pt_size))
            p.font.bold = bold
            p.font.name = 'Microsoft YaHei'
            p.font.color.rgb = text_color
            p.alignment = align

            y_offset += h_estimate + 0.05

    def _render_elements_slide(self, slide, data: Dict):
        """Render a slide directly from persisted Fabric-style elements."""
        for element in data.get("elements", []):
            element_type = element.get("type")
            if element_type == "rect":
                self._add_rect_element(slide, element)
            elif element_type == "circle":
                self._add_circle_element(slide, element)
            elif element_type == "line":
                self._add_line_element(slide, element)
            elif element_type == "image":
                self._add_image_element(slide, element)
            elif element_type == "text":
                self._add_text_element(slide, element)

    def _to_inches_x(self, value: float):
        return Inches(max(0, value * PX_TO_IN_X))

    def _to_inches_y(self, value: float):
        return Inches(max(0, value * PX_TO_IN_Y))

    def _parse_color(self, value: Optional[str], fallback=None):
        if not value:
            return fallback
        if isinstance(value, str):
            text = value.strip()
            if text.startswith("#") and len(text) == 7:
                try:
                    return RGBColor(int(text[1:3], 16), int(text[3:5], 16), int(text[5:7], 16))
                except ValueError:
                    return fallback
        return fallback

    def _theme_color(self, data: Dict, key: str, fallback):
        theme = data.get("theme", {}) if isinstance(data.get("theme"), dict) else {}
        return self._parse_color(theme.get(key), fallback)

    def _add_text_element(self, slide, element: Dict):
        style = element.get("style", {})
        tx_box = slide.shapes.add_textbox(
            self._to_inches_x(element.get("x", 0)),
            self._to_inches_y(element.get("y", 0)),
            self._to_inches_x(element.get("width", 240)),
            self._to_inches_y(element.get("height", 48)),
        )
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(element.get("content", ""))
        font = p.font
        font.size = Pt(style.get("fontSize", 24))
        font.name = style.get("fontFamily", "Arial")
        font.bold = style.get("fontWeight") == "bold"
        font.italic = style.get("fontStyle") == "italic"
        font.underline = bool(style.get("underline"))
        font_color = self._parse_color(style.get("fill"), LIGHT_TEXT)
        if font_color is not None:
            font.color.rgb = font_color

    def _add_rect_element(self, slide, element: Dict):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self._to_inches_x(element.get("x", 0)),
            self._to_inches_y(element.get("y", 0)),
            self._to_inches_x(element.get("width", 100)),
            self._to_inches_y(element.get("height", 100)),
        )
        style = element.get("style", {})
        fill_color = self._parse_color(style.get("fill") or element.get("fill"), None)
        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        stroke_color = self._parse_color(style.get("stroke") or element.get("stroke"), None)
        if stroke_color is not None:
            shape.line.color.rgb = stroke_color
        else:
            shape.line.fill.background()

    def _add_circle_element(self, slide, element: Dict):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            self._to_inches_x(element.get("x", 0)),
            self._to_inches_y(element.get("y", 0)),
            self._to_inches_x(element.get("width", 40)),
            self._to_inches_y(element.get("height", 40)),
        )
        style = element.get("style", {})
        fill_color = self._parse_color(style.get("fill") or element.get("fill"), None)
        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        shape.line.fill.background()

    def _add_line_element(self, slide, element: Dict):
        style = element.get("style", {})
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            self._to_inches_x(element.get("x", 0)),
            self._to_inches_y(element.get("y", 0)),
            self._to_inches_x(element.get("x", 0) + element.get("width", 100)),
            self._to_inches_y(element.get("y", 0) + element.get("height", 0)),
        )
        stroke_color = self._parse_color(style.get("stroke") or element.get("stroke"), LIGHT_TEXT)
        if stroke_color is not None:
            connector.line.color.rgb = stroke_color

    def _add_image_element(self, slide, element: Dict):
        content = element.get("content") or element.get("src")
        if not content:
            return

        image_source = None
        if isinstance(content, str) and content.startswith("data:image"):
            try:
                encoded = content.split(",", 1)[1]
                image_source = BytesIO(base64.b64decode(encoded))
            except Exception:
                return
        elif isinstance(content, str):
            path = Path(content)
            if path.exists():
                image_source = str(path)

        if image_source is None:
            return

        slide.shapes.add_picture(
            image_source,
            self._to_inches_x(element.get("x", 0)),
            self._to_inches_y(element.get("y", 0)),
            width=self._to_inches_x(element.get("width", 200)),
            height=self._to_inches_y(element.get("height", 120)),
        )

    def _add_optional_image(self, slide, data: Dict):
        """Add user-uploaded image if present in slide dict."""
        img_path = data.get("imagePath") or data.get("image_path") or ""
        img_url = data.get("image") or ""

        # Resolve local path from /uploads/... URL when available.
        if not img_path and isinstance(img_url, str) and img_url.startswith("/uploads/"):
            img_path = str((Path.cwd() / img_url.lstrip("/")).resolve())

        if not img_path:
            return

        path_obj = Path(img_path)
        if not path_obj.exists() or not path_obj.is_file():
            return

        # Place image in lower-right area; keep room for source/footer.
        slide.shapes.add_picture(str(path_obj), Inches(7.1), Inches(4.0), width=Inches(5.6), height=Inches(2.9))

    # ── Renderers ─────────────────────────────────────────────

    def _render_title(self, slide, data: Dict):
        title = data.get("title", "")
        subtitle = data.get("subtitle", "")
        title_color = self._theme_color(data, "title", LIGHT_TEXT)
        dim_color = self._theme_color(data, "body", MUTED_TEXT)

        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(1))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(24)
            p2.font.color.rgb = dim_color
            p2.alignment = PP_ALIGN.CENTER

    def _render_bullets(self, slide, data: Dict):
        title = data.get("title", "")
        items = data.get("items", [])
        title_color = self._theme_color(data, "title", LIGHT_TEXT)
        body_color = self._theme_color(data, "body", RGBColor(0xC9, 0xD1, 0xD9))

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = title_color
        p.font.bold = True

        # Bullets
        y = 2.0
        for item in items[:6]:
            text = item if isinstance(item, str) else str(item)
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(10), Inches(0.7))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {text}"
            p.font.size = Pt(20)
            p.font.color.rgb = body_color
            y += 0.85

    def _render_comparison(self, slide, data: Dict):
        title = data.get("title", "")
        left = data.get("left", {})
        right = data.get("right", {})
        title_color = self._theme_color(data, "title", LIGHT_TEXT)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = title_color
        p.font.bold = True

        # Left column
        self._add_compare_column(slide, left, Inches(0.8), "←")
        # Right column
        self._add_compare_column(slide, right, Inches(6.8), "→")

    def _add_compare_column(self, slide, col, left_offset, arrow: str):
        # Handle string columns from LLM
        if isinstance(col, str):
            col = {"label": col, "desc": "", "items": []}
        if not isinstance(col, dict):
            col = {"label": str(col), "desc": "", "items": []}
        label = col.get("label", "")
        desc = col.get("desc", "")
        items = col.get("items", [])

        txBox = slide.shapes.add_textbox(left_offset, Inches(2.0), Inches(5.5), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{arrow} {label}"
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_BLUE
        p.font.bold = True

        text = desc
        if items:
            text += "\n" + "\n".join(f"• {i}" for i in items[:4])

        txBox2 = slide.shapes.add_textbox(left_offset, Inches(2.8), Inches(5.5), Inches(3.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = text
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xC9, 0xD1, 0xD9)

    def _render_metric(self, slide, data: Dict):
        title = data.get("title", "")
        number = str(data.get("number", ""))
        label = data.get("label", "")
        items = data.get("items", [])
        accent_color = self._theme_color(data, "accent", ACCENT_BLUE)
        dim_color = self._theme_color(data, "body", MUTED_TEXT)
        body_color = self._theme_color(data, "body", RGBColor(0xC9, 0xD1, 0xD9))

        # Big number
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(72)
        p.font.color.rgb = accent_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Label
        if label:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(0.8))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = label
            p2.font.size = Pt(24)
            p2.font.color.rgb = dim_color
            p2.alignment = PP_ALIGN.CENTER

        # Items
        y = 4.5
        for item in items[:4]:
            text = item if isinstance(item, str) else str(item)
            txBox = slide.shapes.add_textbox(Inches(2), Inches(y), Inches(9), Inches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {text}"
            p.font.size = Pt(16)
            p.font.color.rgb = body_color
            y += 0.6

    def _render_quote(self, slide, data: Dict):
        quote_text = data.get("quote", "")
        attribution = data.get("attribution", "")

        # Quote marks
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "\u201C"
        p.font.size = Pt(80)
        p.font.color.rgb = ACCENT_BLUE

        # Quote text
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = quote_text
        p2.font.size = Pt(28)
        p2.font.color.rgb = LIGHT_TEXT
        p2.font.italic = True

        # Attribution
        if attribution:
            txBox3 = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7), Inches(0.6))
            tf3 = txBox3.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = f"— {attribution}"
            p3.font.size = Pt(18)
            p3.font.color.rgb = MUTED_TEXT
            p3.alignment = PP_ALIGN.RIGHT

    def _render_timeline(self, slide, data: Dict):
        title = data.get("title", "")
        events = data.get("events", [])

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = LIGHT_TEXT
        p.font.bold = True

        y = 2.0
        for evt in events[:5]:
            year = evt.get("year", "")
            text = evt.get("text", "")

            # Year marker
            txBox = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(2), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(year)
            p.font.size = Pt(20)
            p.font.color.rgb = ACCENT_BLUE
            p.font.bold = True

            # Event text
            txBox2 = slide.shapes.add_textbox(Inches(3.5), Inches(y), Inches(8), Inches(0.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = text
            p2.font.size = Pt(16)
            p2.font.color.rgb = LIGHT_TEXT

            y += 1.0

    def _render_flowchart(self, slide, data: Dict):
        title = data.get("title", "")
        steps = data.get("steps", [])

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = LIGHT_TEXT
        p.font.bold = True

        if not steps:
            return

        # Calculate positions for steps
        num_steps = min(len(steps), 6)
        step_width = 1.8
        total_width = num_steps * step_width + (num_steps - 1) * 0.5
        start_x = (13.333 - total_width) / 2

        for i, step in enumerate(steps[:num_steps]):
            x = start_x + i * (step_width + 0.5)
            label = step if isinstance(step, str) else step.get("label", str(step))

            # Step box
            txBox = slide.shapes.add_textbox(Inches(x), Inches(3.5), Inches(step_width), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = LIGHT_TEXT
            p.alignment = PP_ALIGN.CENTER

            # Arrow between steps
            if i < num_steps - 1:
                arrow_x = x + step_width
                txBox_arrow = slide.shapes.add_textbox(Inches(arrow_x), Inches(3.7), Inches(0.5), Inches(0.5))
                tf_arrow = txBox_arrow.text_frame
                p_arrow = tf_arrow.paragraphs[0]
                p_arrow.text = "→"
                p_arrow.font.size = Pt(20)
                p_arrow.font.color.rgb = ACCENT_BLUE
                p_arrow.alignment = PP_ALIGN.CENTER

    # ── Helpers ───────────────────────────────────────────────

    def _add_source(self, slide, source: str):
        """Add source annotation at bottom-left."""
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"来源: {source}"
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED_TEXT

    # ── Template-based export ──────────────────────────────────

    def _export_with_template(
        self,
        slides: List[Dict[str, Any]],
        output_path: str,
        template_path: str,
    ) -> str:
        """Export using a .pptx template — fill placeholder layouts with content."""
        from backend.ppt.template_manager import TemplateManager

        prs = PptxPresentation(template_path)

        # Remove existing slides from template
        while len(prs.slides) > 0:
            from pptx.oxml.ns import qn
            rId = prs.slides._sldIdLst[0].get(qn("r:id"))
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

        # Build layout lookup
        layout_map = {}
        for i, layout in enumerate(prs.slide_layouts):
            layout_map[i] = layout
            layout_map[layout.name.lower()] = layout

        for slide_data in slides:
            layout = self._resolve_layout(prs, layout_map, slide_data)
            slide = prs.slides.add_slide(layout)
            self._fill_placeholders(slide, slide_data)

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        return str(out.absolute())

    def _resolve_layout(self, prs, layout_map, slide_data):
        """Resolve which layout to use for a slide."""
        # Explicit layout_index
        idx = slide_data.get("layout_index")
        if idx is not None and idx in layout_map:
            return layout_map[idx]

        # Explicit layout_name (partial match)
        name = slide_data.get("layout_name", "").lower()
        if name:
            for key, layout in layout_map.items():
                if isinstance(key, str) and name in key:
                    return layout

        # Infer from slide type
        slide_type = slide_data.get("type", "bullets")
        type_layout_hints = {
            "title": ["title", "标题", "cover", "封面"],
            "section": ["section", "章节", "divider"],
            "bullets": ["content", "内容", "body", "text", "正文"],
            "comparison": ["two content", "对比", "comparison", "两栏"],
            "metric": ["blank", "空白"],
            "quote": ["blank", "空白"],
            "timeline": ["content", "内容"],
            "flowchart": ["content", "内容"],
            "ending": ["title", "标题", "cover", "封面", "ending", "结尾"],
        }
        hints = type_layout_hints.get(slide_type, ["content", "blank"])
        for hint in hints:
            for key, layout in layout_map.items():
                if isinstance(key, str) and hint in key:
                    return layout

        # Fallback: first layout with 2+ placeholders, else layout 0
        for layout in prs.slide_layouts:
            if len(list(layout.placeholders)) >= 2:
                return layout
        return prs.slide_layouts[0]

    def _fill_placeholders(self, slide, slide_data: Dict[str, Any]):
        """Fill slide placeholders from structured data."""
        # Explicit placeholder mapping
        ph_map = slide_data.get("placeholders", {})

        # Auto-map common fields to standard placeholder indices
        if "title" in slide_data:
            ph_map.setdefault(0, slide_data["title"])
            ph_map.setdefault("0", slide_data["title"])
        if "subtitle" in slide_data:
            ph_map.setdefault(1, slide_data["subtitle"])
        if "body" in slide_data:
            ph_map.setdefault(1, slide_data["body"])
        if "items" in slide_data and isinstance(slide_data["items"], list):
            # Format bullet items as body text
            body_text = "\n".join(f"• {item}" for item in slide_data["items"])
            ph_map.setdefault(1, body_text)

        # Build compound body for complex types
        slide_type = slide_data.get("type", "")
        if slide_type == "comparison":
            left = slide_data.get("left", {})
            right = slide_data.get("right", {})
            if isinstance(left, dict) and isinstance(right, dict):
                left_text = self._format_compare_column(left)
                right_text = self._format_compare_column(right)
                ph_map.setdefault(1, left_text)
                ph_map.setdefault(2, right_text)
        elif slide_type == "metric":
            number = slide_data.get("number", "")
            label = slide_data.get("label", "")
            items = slide_data.get("items", [])
            body = f"{number}\n{label}"
            if items:
                body += "\n" + "\n".join(f"• {i}" for i in items)
            ph_map.setdefault(1, body)
        elif slide_type == "quote":
            quote = slide_data.get("quote", "")
            attribution = slide_data.get("attribution", "")
            ph_map.setdefault(1, f"\"{quote}\"\n\n— {attribution}" if attribution else f"\"{quote}\"")
        elif slide_type == "timeline":
            events = slide_data.get("events", [])
            lines = [f"{e.get('year', '')}: {e.get('text', '')}" for e in events]
            ph_map.setdefault(1, "\n".join(lines))
        elif slide_type == "flowchart":
            steps = slide_data.get("steps", [])
            step_text = " → ".join(str(s) if isinstance(s, str) else s.get("label", str(s)) for s in steps)
            ph_map.setdefault(1, step_text)

        # Fill each placeholder
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            content = ph_map.get(idx) or ph_map.get(str(idx))
            if content is None:
                continue

            if isinstance(content, str):
                ph.text = content
            elif isinstance(content, list):
                tf = ph.text_frame
                tf.clear()
                for i, item in enumerate(content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = str(item)

    def _format_compare_column(self, col) -> str:
        """Format a comparison column dict as text."""
        if isinstance(col, str):
            return col
        label = col.get("label", "")
        desc = col.get("desc", "")
        items = col.get("items", [])
        parts = []
        if label:
            parts.append(label)
        if desc:
            parts.append(desc)
        if items:
            parts.extend(f"• {i}" for i in items)
        return "\n".join(parts)
