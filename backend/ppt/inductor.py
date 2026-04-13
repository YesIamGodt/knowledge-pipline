"""
PPTAgent SlideInducter 集成
从 PPT 模板中提取布局模式
"""
import os
from typing import List
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches

from backend.ppt.models import SlideLayout, LayoutElement


class SlideInducter:
    """
    从 PPT 模板中提取布局模式

    基于 PPTAgent 的 SlideInducter 设计，简化实现：
    1. 读取 PPT 幻灯片
    2. 分析每页的形状位置和类型
    3. 聚类相似的布局
    4. 生成布局 schema
    """

    # Font size thresholds for text classification
    TITLE_FONT_SIZE = 32
    SUBTITLE_FONT_SIZE = 20

    # Text length thresholds
    TITLE_TEXT_LENGTH = 20
    SUBTITLE_TEXT_LENGTH = 50

    def __init__(self):
        self.layouts: List[SlideLayout] = []

    def analyze(self, pptx_path: str) -> List[SlideLayout]:
        """
        分析 PPT 模板，提取布局模式

        Args:
            pptx_path: PPT 文件路径

        Returns:
            布局模式列表
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"Template file not found: {pptx_path}")

        try:
            prs = Presentation(pptx_path)

            # 获取幻灯片尺寸
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # 提取每页的元素
            slides_data = []
            for slide_idx, slide in enumerate(prs.slides):
                elements = self._extract_slide_elements(slide, slide_idx, slide_width, slide_height)
                slides_data.append({
                    'index': slide_idx,
                    'elements': elements
                })

            # 聚类相似的布局
            layout_clusters = self._cluster_similar_layouts(slides_data)

            # 为每个聚类生成 layout schema
            self.layouts = []
            for cluster_name, cluster_data in layout_clusters.items():
                layout = self._generate_layout_schema(cluster_name, cluster_data)
                self.layouts.append(layout)

            return self.layouts

        except Exception as e:
            raise RuntimeError(f"Failed to analyze template: {str(e)}")

    def _extract_slide_elements(self, slide, slide_idx: int, slide_width, slide_height) -> List[LayoutElement]:
        """提取幻灯片中的元素"""
        elements = []

        # Handle empty slides
        if not slide.shapes:
            return elements

        for shape in slide.shapes:
            # Skip zero-size shapes to prevent division by zero
            if shape.width == 0 or shape.height == 0:
                continue
            if slide_width == 0 or slide_height == 0:
                continue
            # 跳过隐藏的形状（使用 hasattr 检查，因为不是所有形状都有 visible 属性）
            if hasattr(shape, 'visible') and not shape.visible:
                continue

            # 获取位置（归一化到 0-1）
            x = shape.left / slide_width
            y = shape.top / slide_height
            width = shape.width / slide_width
            height = shape.height / slide_height

            # 判断类型
            if hasattr(shape, 'text') and shape.text.strip():
                # 文本框
                element_type = self._classify_text_element(shape)
                elements.append(LayoutElement(
                    type=element_type,
                    x=x,
                    y=y,
                    width=width,
                    height=height,
                    align=self._detect_alignment(shape)
                ))

            elif shape.shape_type == 13:  # Picture
                elements.append(LayoutElement(
                    type='image',
                    x=x,
                    y=y,
                    width=width,
                    height=height
                ))

        return elements

    def _classify_text_element(self, shape) -> str:
        """判断文本元素的类型"""
        text = shape.text.strip()

        # 根据字体大小判断
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font_size = run.font.size
                    if font_size:
                        # 转换为磅
                        size_pt = font_size.pt if hasattr(font_size, 'pt') else font_size

                        if size_pt > self.TITLE_FONT_SIZE:
                            return 'title'
                        elif size_pt > self.SUBTITLE_FONT_SIZE:
                            return 'subtitle'
                        else:
                            return 'text'

        # 根据文本长度判断
        if len(text) < self.TITLE_TEXT_LENGTH and '\n' not in text:
            return 'title'
        elif len(text) < self.SUBTITLE_TEXT_LENGTH:
            return 'subtitle'
        else:
            return 'text'

    def _detect_alignment(self, shape) -> str:
        """检测文本对齐方式"""
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                alignment = paragraph.alignment
                if alignment == 0:  # Left
                    return 'left'
                elif alignment == 1:  # Center
                    return 'center'
                elif alignment == 2:  # Right
                    return 'right'
                elif alignment == 3:  # Justify
                    return 'left'
        return 'left'

    def _cluster_similar_layouts(self, slides_data: List[dict]) -> dict:
        """聚类相似的布局"""
        # 简化实现：根据元素数量和位置分组
        clusters = defaultdict(list)

        for slide_data in slides_data:
            elements = slide_data['elements']

            # 生成布局签名（元素类型 + 位置区间）
            signature = self._generate_layout_signature(elements)

            clusters[signature].append(slide_data)

        # 为每个聚类生成名称
        layout_clusters = {}
        for signature, cluster_slides in clusters.items():
            layout_name = self._generate_layout_name(signature, cluster_slides[0]['elements'])
            layout_clusters[layout_name] = {
                'slides': cluster_slides,
                'sample_slide': cluster_slides[0]
            }

        return layout_clusters

    def _generate_layout_signature(self, elements: List[LayoutElement]) -> str:
        """生成布局签名"""
        parts = []

        # 按元素类型分组
        type_counts = defaultdict(int)
        for elem in elements:
            type_counts[elem.type] += 1

        # 生成签名
        if type_counts['title'] > 0:
            parts.append('title')
        if type_counts['subtitle'] > 0:
            parts.append('subtitle')
        if type_counts['text'] > 0:
            parts.append(f"text{type_counts['text']}")
        if type_counts['image'] > 0:
            parts.append(f"img{type_counts['image']}")

        return '-'.join(parts) if parts else 'blank'

    def _generate_layout_name(self, signature: str, elements: List[LayoutElement]) -> str:
        """生成布局名称"""
        # 根据签名和元素位置生成描述性名称
        if 'title' in signature and 'subtitle' in signature:
            return 'title-center'
        elif 'text2' in signature:
            return 'two-column-text'
        elif 'img1' in signature and 'text1' in signature:
            return 'image-with-text'
        else:
            return signature.replace('-', '-')

    def _generate_layout_schema(self, name: str, cluster_data: dict) -> SlideLayout:
        """生成布局 schema"""
        sample_elements = cluster_data['sample_slide']['elements']

        # 生成描述
        description = self._generate_layout_description(name, sample_elements)

        return SlideLayout(
            name=name,
            description=description,
            elements=sample_elements,
            sample_slide_index=cluster_data['sample_slide']['index']
        )

    def _generate_layout_description(self, name: str, elements: List[LayoutElement]) -> str:
        """生成布局描述"""
        if name == 'title-center':
            return '标题居中，副标题下方'
        elif name == 'two-column-text':
            return '双栏文字布局'
        elif name == 'image-with-text':
            return '图片配文字说明'
        else:
            return f'自定义布局: {name}'
