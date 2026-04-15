"""
PPTX Template Manager — upload, parse, store, and apply user templates.

Core flow:
  1. User uploads .pptx → parse_template() extracts layouts + placeholders
  2. LLM generates structured slide JSON
  3. generate_from_template() fills template layouts with LLM content → native PPTX
"""

import json
import shutil
from copy import deepcopy
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# Where uploaded templates are stored
TEMPLATE_DIR = Path(__file__).resolve().parents[2] / "uploads" / "templates"


class LayoutInfo:
    """Parsed info about a single slide layout from a template."""

    def __init__(self, index: int, name: str, placeholders: List[Dict[str, Any]]):
        self.index = index
        self.name = name
        self.placeholders = placeholders  # [{idx, type, name, left, top, width, height}]

    def to_dict(self) -> Dict[str, Any]:
        return {
            "index": self.index,
            "name": self.name,
            "placeholders": self.placeholders,
        }


class TemplateInfo:
    """Full parsed template information."""

    def __init__(self, template_id: str, filename: str, path: str):
        self.template_id = template_id
        self.filename = filename
        self.path = path
        self.layouts: List[LayoutInfo] = []
        self.colors: List[str] = []
        self.fonts: List[str] = []
        self.slide_width: float = 0
        self.slide_height: float = 0

    def to_dict(self) -> Dict[str, Any]:
        return {
            "template_id": self.template_id,
            "filename": self.filename,
            "path": self.path,
            "layouts": [l.to_dict() for l in self.layouts],
            "colors": self.colors,
            "fonts": self.fonts,
            "slide_width": self.slide_width,
            "slide_height": self.slide_height,
        }

    def get_layout_by_name(self, name: str) -> Optional[LayoutInfo]:
        """Find layout by name (case-insensitive partial match)."""
        name_lower = name.lower()
        for layout in self.layouts:
            if name_lower in layout.name.lower():
                return layout
        return None

    def get_layout_by_index(self, index: int) -> Optional[LayoutInfo]:
        """Find layout by index."""
        for layout in self.layouts:
            if layout.index == index:
                return layout
        return None

    def to_prompt_text(self) -> str:
        """Generate LLM-readable description of available layouts."""
        lines = [f"模板: {self.filename}", "可用布局:"]
        for layout in self.layouts:
            ph_desc = []
            for ph in layout.placeholders:
                ph_desc.append(f"{ph['name']}(idx={ph['idx']})")
            phs = ", ".join(ph_desc) if ph_desc else "无占位符"
            lines.append(f"  - layout_index={layout.index}: \"{layout.name}\" → 占位符: {phs}")
        return "\n".join(lines)


class TemplateManager:
    """Manages uploaded PPTX templates."""

    # Map placeholder type ids to human-readable names
    PH_TYPE_MAP = {
        0: "title",
        1: "body",
        2: "center_title",
        3: "subtitle",
        4: "date",
        5: "footer",
        6: "slide_number",
        7: "object",
        10: "picture",
        12: "table",
        13: "org_chart",
        14: "media_clip",
        15: "title",  # TITLE in some templates
    }

    def __init__(self):
        if not HAS_PPTX:
            raise ImportError("python-pptx required: pip install python-pptx")
        TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)
        self._meta_file = TEMPLATE_DIR / "_templates.json"
        self._templates: Dict[str, TemplateInfo] = {}
        self._load_meta()

    def _load_meta(self):
        """Load template metadata from disk."""
        if self._meta_file.exists():
            try:
                data = json.loads(self._meta_file.read_text(encoding="utf-8"))
                for tid, info in data.items():
                    tpl = TemplateInfo(tid, info["filename"], info["path"])
                    tpl.colors = info.get("colors", [])
                    tpl.fonts = info.get("fonts", [])
                    tpl.slide_width = info.get("slide_width", 0)
                    tpl.slide_height = info.get("slide_height", 0)
                    for l in info.get("layouts", []):
                        tpl.layouts.append(LayoutInfo(l["index"], l["name"], l["placeholders"]))
                    self._templates[tid] = tpl
            except (json.JSONDecodeError, KeyError):
                pass

    def _save_meta(self):
        """Persist template metadata."""
        data = {tid: tpl.to_dict() for tid, tpl in self._templates.items()}
        self._meta_file.write_text(
            json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
        )

    def list_templates(self) -> List[Dict[str, Any]]:
        """List all available templates."""
        result = []
        for tid, tpl in self._templates.items():
            result.append({
                "template_id": tid,
                "filename": tpl.filename,
                "layouts": len(tpl.layouts),
                "colors": tpl.colors[:4],
                "fonts": tpl.fonts[:3],
            })
        return result

    def get_template(self, template_id: str) -> Optional[TemplateInfo]:
        """Get full template info by ID."""
        return self._templates.get(template_id)

    def upload_template(self, source_path: str, template_id: Optional[str] = None) -> TemplateInfo:
        """
        Upload and parse a .pptx template.

        Args:
            source_path: Path to the .pptx file
            template_id: Optional custom ID (auto-generated from filename if not provided)

        Returns:
            Parsed TemplateInfo
        """
        src = Path(source_path)
        if not src.exists():
            raise FileNotFoundError(f"Template not found: {source_path}")
        if not src.suffix.lower() == ".pptx":
            raise ValueError("Only .pptx files are supported")

        # Generate template ID
        if not template_id:
            template_id = src.stem.lower().replace(" ", "-").replace("_", "-")
            # Avoid collision
            if template_id in self._templates:
                i = 2
                while f"{template_id}-{i}" in self._templates:
                    i += 1
                template_id = f"{template_id}-{i}"

        # Copy to templates directory
        dest = TEMPLATE_DIR / f"{template_id}.pptx"
        shutil.copy2(str(src), str(dest))

        # Parse the template
        tpl = self._parse_template(template_id, src.name, str(dest))
        self._templates[template_id] = tpl
        self._save_meta()
        return tpl

    def delete_template(self, template_id: str) -> bool:
        """Delete an uploaded template."""
        if template_id not in self._templates:
            return False
        tpl = self._templates[template_id]
        try:
            Path(tpl.path).unlink(missing_ok=True)
        except OSError:
            pass
        del self._templates[template_id]
        self._save_meta()
        return True

    def _parse_template(self, template_id: str, filename: str, path: str) -> TemplateInfo:
        """Parse a .pptx file to extract layout and style information."""
        tpl = TemplateInfo(template_id, filename, path)
        prs = PptxPresentation(path)

        # Slide dimensions
        tpl.slide_width = prs.slide_width.inches if prs.slide_width else 13.333
        tpl.slide_height = prs.slide_height.inches if prs.slide_height else 7.5

        # Extract layouts
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                ph_info = {
                    "idx": ph.placeholder_format.idx,
                    "type": self.PH_TYPE_MAP.get(ph.placeholder_format.type, "unknown"),
                    "name": ph.name or f"Placeholder {ph.placeholder_format.idx}",
                    "left": round(ph.left.inches, 3) if ph.left else 0,
                    "top": round(ph.top.inches, 3) if ph.top else 0,
                    "width": round(ph.width.inches, 3) if ph.width else 0,
                    "height": round(ph.height.inches, 3) if ph.height else 0,
                }
                placeholders.append(ph_info)
            tpl.layouts.append(LayoutInfo(i, layout.name, placeholders))

        # Extract theme colors (reuse logic from template_analyzer)
        tpl.colors = self._extract_colors(prs)
        tpl.fonts = self._extract_fonts(prs)

        return tpl

    def _extract_colors(self, prs: 'PptxPresentation') -> List[str]:
        """Extract theme colors from presentation."""
        colors = []
        try:
            theme = prs.slide_masters[0].element.find(
                ".//" + qn("a:clrScheme")
            )
            if theme is not None:
                for child in theme:
                    for sub in child:
                        val = sub.get("val") or sub.get("lastClr")
                        if val and len(val) == 6:
                            c = f"#{val}"
                            if c not in colors:
                                colors.append(c)
        except Exception:
            pass
        return colors

    def _extract_fonts(self, prs: 'PptxPresentation') -> List[str]:
        """Extract theme fonts from presentation."""
        fonts = []
        try:
            theme_elem = prs.slide_masters[0].element
            for tag in ("latin", "ea", "cs"):
                for elem in theme_elem.iter(qn(f"a:{tag}")):
                    typeface = elem.get("typeface")
                    if typeface and typeface not in fonts and not typeface.startswith("+"):
                        fonts.append(typeface)
        except Exception:
            pass
        return fonts

    def generate_from_template(
        self,
        template_id: str,
        slides_data: List[Dict[str, Any]],
        output_path: str,
    ) -> str:
        """
        Generate a PPTX file by filling template layouts with content.

        Each slide_data dict should have:
          - layout_index: int (which layout to use)
          - OR layout_name: str (partial match)
          - placeholders: dict mapping placeholder idx → content string
          - OR title/body/subtitle: shorthand for common placeholders

        Args:
            template_id: ID of uploaded template
            slides_data: List of slide content dicts
            output_path: Where to save the generated .pptx

        Returns:
            Absolute path of generated file
        """
        tpl = self._templates.get(template_id)
        if not tpl:
            raise ValueError(f"Template not found: {template_id}")
        if not Path(tpl.path).exists():
            raise FileNotFoundError(f"Template file missing: {tpl.path}")

        prs = PptxPresentation(tpl.path)

        # Remove any existing slides from the template (keep only masters/layouts)
        self._remove_existing_slides(prs)

        # Generate slides from data
        for slide_data in slides_data:
            self._add_slide_from_data(prs, tpl, slide_data)

        # Save
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        return str(out.absolute())

    def _remove_existing_slides(self, prs: 'PptxPresentation'):
        """Remove all slides from a presentation (keep layouts)."""
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(qn("r:id"))
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    def _add_slide_from_data(
        self,
        prs: 'PptxPresentation',
        tpl: TemplateInfo,
        slide_data: Dict[str, Any],
    ):
        """Add a single slide using a template layout and fill placeholders."""
        # Determine which layout to use
        layout_index = slide_data.get("layout_index")
        layout_name = slide_data.get("layout_name", "")

        layout = None
        if layout_index is not None and 0 <= layout_index < len(prs.slide_layouts):
            layout = prs.slide_layouts[layout_index]
        elif layout_name:
            # Find by name (partial match)
            for i, sl in enumerate(prs.slide_layouts):
                if layout_name.lower() in sl.name.lower():
                    layout = sl
                    break

        if layout is None:
            # Fallback: use first layout with placeholders, or layout 0
            for sl in prs.slide_layouts:
                if len(list(sl.placeholders)) > 1:
                    layout = sl
                    break
            if layout is None:
                layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        # Fill placeholders
        ph_map = slide_data.get("placeholders", {})

        # Also support shorthand: title → idx 0, body → idx 1, subtitle → idx 1
        if "title" in slide_data and 0 not in ph_map and "0" not in ph_map:
            ph_map[0] = slide_data["title"]
        if "body" in slide_data:
            # Body goes to idx 1 (most common)
            if 1 not in ph_map and "1" not in ph_map:
                ph_map[1] = slide_data["body"]
        if "subtitle" in slide_data:
            if 1 not in ph_map and "1" not in ph_map:
                ph_map[1] = slide_data["subtitle"]

        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            # Try int key first, then string key
            content = ph_map.get(idx) or ph_map.get(str(idx))
            if content is None:
                continue

            if isinstance(content, str):
                ph.text = content
            elif isinstance(content, list):
                # List of bullet items
                tf = ph.text_frame
                tf.clear()
                for i, item in enumerate(content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0
            elif isinstance(content, dict):
                # Rich content with formatting
                self._fill_rich_placeholder(ph, content)

    def _fill_rich_placeholder(self, ph, content: Dict[str, Any]):
        """Fill a placeholder with rich content (text + formatting)."""
        text = content.get("text", "")
        items = content.get("items", [])
        font_size = content.get("font_size")
        bold = content.get("bold")
        color = content.get("color")

        tf = ph.text_frame
        tf.clear()

        if items:
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(item)
                p.level = 0
        elif text:
            tf.paragraphs[0].text = text

        # Apply formatting
        for p in tf.paragraphs:
            for run in p.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                if color and isinstance(color, str) and color.startswith("#") and len(color) == 7:
                    try:
                        run.font.color.rgb = RGBColor(
                            int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
                        )
                    except ValueError:
                        pass
