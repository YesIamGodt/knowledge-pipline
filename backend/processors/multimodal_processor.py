"""
多模态文件处理器 - 使用 Claude Code 内置多模态能力处理图片

核心策略：
1. 文字提取 → Python (PyMuPDF, python-docx 等) 毫秒级完成
2. 图片理解 → Claude Code Read 工具原生理解，无需外部 API
3. Python 只负责提取图片字节，由 Claude Code agent 处理图片理解

注意：图片不通过外部 LLM API 处理，而是由 Claude Code 的内置多模态工具处理。
"""

import logging
import base64
import io
import hashlib
import threading
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from .file_processor import ProcessedDocument

logger = logging.getLogger(__name__)


class MultimodalProcessor:
    """
    多模态大模型增强的文件处理器

    使用已配置的多模态大模型（如 MiniMax-M2.7）来：
    - 直接理解图片内容（无需 OCR）
    - 处理复杂文档格式
    - 智能提取表格、图表、公式
    - 处理 PDF 中的图片和扫描内容
    - 处理 PPT 中的图片内容
    """

    def __init__(self, config: Optional[Dict[str, str]] = None):
        """
        初始化多模态处理器

        Args:
            config: LLM 配置字典，包含 base_url, model, api_key
        """
        self.config = config or {}
        self.client = None
        self.model = self.config.get("model", "MiniMax-M2.7")
        self._progress_callback = None  # 进度回调
        self._init_client()

    def set_progress_callback(self, callback):
        """设置进度回调函数: callback(message: str)"""
        self._progress_callback = callback

    def _emit_progress(self, message: str):
        """发送进度信息到终端"""
        if self._progress_callback:
            self._progress_callback(message)
        else:
            import sys
            print(f"  {message}", flush=True, file=sys.stderr)

    def _init_client(self):
        """初始化 OpenAI 兼容客户端"""
        try:
            import openai
            self.client = openai.OpenAI(
                base_url=self.config.get("base_url"),
                api_key=self.config.get("api_key")
            )
            logger.info(f"多模态处理器初始化成功，使用模型: {self.model}")
        except ImportError:
            logger.warning("openai 库未安装，多模态处理不可用")
        except Exception as e:
            logger.warning(f"多模态客户端初始化失败: {e}")

    def is_available(self) -> bool:
        """检查多模态处理是否可用"""
        return self.client is not None

    def process_file(self, file_path: str) -> ProcessedDocument:
        """
        处理文件，根据文件类型选择最佳方法

        优先使用多模态大模型处理复杂场景，
        对于简单文本文件，仍然可以使用传统方法。
        """
        path = Path(file_path)

        if not path.exists():
            return ProcessedDocument(
                content="",
                errors=[f"文件不存在: {file_path}"]
            )

        ext = path.suffix.lower()

        # 判断是否需要多模态处理
        needs_multimodal = self._needs_multimodal(ext, path)

        if needs_multimodal and self.is_available():
            return self._process_with_multimodal(path, ext)
        else:
            # 回退到传统方法或简单处理
            return self._process_fallback(path, ext)

    def _needs_multimodal(self, ext: str, path: Path) -> bool:
        """判断文件是否需要多模态处理"""
        multimodal_extensions = {
            ".jpg", ".jpeg", ".png", ".webp",  # 图片
            ".pdf",  # PDF（可能包含图片或扫描内容）
            ".bmp", ".gif", ".tiff",  # 其他图片格式
            ".pptx",  # PPT（可能包含图片）
            ".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv",  # 视频
            ".webm", ".m4v", ".mpg", ".mpeg", ".3gp", ".ts",  # 更多视频格式
        }
        return ext in multimodal_extensions

    def _process_with_multimodal(self, path: Path, ext: str) -> ProcessedDocument:
        """使用多模态大模型处理文件"""
        errors = []
        content_parts = []
        metadata = {
            "source": str(path),
            "type": ext.lstrip('.'),
            "processor": "multimodal_llm",
            "model": self.model,
        }

        try:
            if ext in [".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"]:
                # 处理图片
                content = self._process_image_with_llm(path)
                content_parts.append(content)

            elif ext == ".pdf":
                # 处理 PDF（文字 + 图片）
                pdf_result = self._process_pdf_multimodal(path)
                content_parts.extend(pdf_result.get("content_parts", []))
                if pdf_result.get("images_processed", 0) > 0:
                    metadata["images_processed"] = pdf_result["images_processed"]

            elif ext == ".pptx":
                # 处理 PPTX（文字 + 图片）
                pptx_result = self._process_pptx_multimodal(path)
                content_parts.extend(pptx_result.get("content_parts", []))
                if pptx_result.get("images_processed", 0) > 0:
                    metadata["images_processed"] = pptx_result["images_processed"]

            elif ext in {".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv",
                         ".webm", ".m4v", ".mpg", ".mpeg", ".3gp", ".ts"}:
                # 处理视频（关键帧提取 + 多模态理解）
                video_result = self._process_video_multimodal(path)
                content_parts.extend(video_result.get("content_parts", []))
                if video_result.get("frames_processed", 0) > 0:
                    metadata["frames_processed"] = video_result["frames_processed"]
                if video_result.get("duration_sec"):
                    metadata["duration_sec"] = video_result["duration_sec"]
                if video_result.get("has_audio_transcript"):
                    metadata["has_audio_transcript"] = True

            else:
                # 其他格式，尝试读取后让模型处理
                content = self._process_other_with_llm(path, ext)
                content_parts.append(content)

            content = "\n\n".join(filter(None, content_parts))
            metadata["processing_status"] = "success"

        except Exception as e:
            logger.error(f"多模态处理失败: {e}")
            errors.append(f"多模态处理失败: {str(e)}")
            metadata["processing_status"] = "failed"
            # 回退到简单处理
            fallback = self._process_fallback(path, ext)
            content = fallback.content
            errors.extend(fallback.errors)

        return ProcessedDocument(
            content=content,
            metadata=metadata,
            errors=errors
        )

    # 多模态处理限制参数
    MAX_IMAGES_PER_DOC = 24       # 每份文档最多处理的图片数
    MIN_IMAGE_SIZE = 5000         # 最小图片字节数，过小的通常是 logo/图标
    MAX_TOKENS_PER_IMAGE = 512    # 单张图片 LLM 响应最多 token 数（防止过长）
    LLM_TIMEOUT = 120              # 单次 LLM 调用超时（秒），kimi-k2.5 处理压缩图片约需 60s
    MAX_PDF_PAGE_RENDER = 5       # PDF 页面渲染最多处理页数
    PAGE_TEXT_CONTEXT_CHARS = 1200  # 传给图片理解的页内文本上下文长度
    MIN_IMAGE_RETAIN = 8          # 多图文档的最少保留图片数

    # 图片压缩参数（优化版 - 2026-04-10）
    # 实测数据：PDF 中 10/20 张图片 >= 50KB，优化后 1278KB → 409KB（节省 68%）
    COMPRESS_THRESHOLD_KB = 50    # 触发压缩的最小图片大小（原 200KB 太高，导致 0 图片被压缩）
    COMPRESS_MAX_DIM = 600        # resize 最大尺寸（原 800px，600px 足够内容理解）
    COMPRESS_TARGET_KB = 80       # 压缩目标大小（原 150KB）
    COMPRESS_QUALITIES = [70, 50, 30]  # 渐进质量（原 [80, 60, 40]）

    def _process_pdf_multimodal(self, pdf_path: Path) -> Dict[str, Any]:
        """
        处理 PDF - 同时提取文字和图片，用多模态大模型理解图片

        Returns:
            包含 content_parts 和 images_processed 的字典
        """
        result = {
            "content_parts": [],
            "images_processed": 0
        }

        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF 未安装，无法进行 PDF 多模态处理")
            return result

        try:
            doc = fitz.open(str(pdf_path))
            logger.info(f"开始处理 PDF，共 {len(doc)} 页")

            # 1. 首先尝试提取文字
            text_content = self._extract_pdf_text(doc)
            if text_content.strip():
                result["content_parts"].append("=== PDF 文字内容 ===\n" + text_content)
                logger.info(f"提取到 PDF 文字: {len(text_content)} 字符")

            # 2. 构建页内文本索引，为图片理解提供上下文
            page_text_map = self._extract_pdf_text_by_page(doc)

            # 3. 提取并处理 PDF 中的图片（去重 + 评分 + 上下文融合）
            images_result = self._extract_and_process_pdf_images(doc, pdf_path, page_text_map)
            if images_result["content"].strip():
                result["content_parts"].append("=== PDF 图片内容 ===\n" + images_result["content"])
            result["images_processed"] = images_result["count"]
            result["images_found"] = images_result.get("found", 0)
            result["images_selected"] = images_result.get("selected", 0)

            # 4. 如果文字提取很少或没有，尝试将页面渲染为图片处理（优先文本稀疏页）
            if len(text_content.strip()) < 500:
                logger.info("PDF 文字较少，可能包含扫描内容，尝试渲染页面为图片处理")
                page_images_content = self._render_and_process_pdf_pages(doc, page_text_map=page_text_map)
                if page_images_content.strip():
                    result["content_parts"].append("=== PDF 页面渲染内容 ===\n" + page_images_content)

            doc.close()

        except Exception as e:
            logger.error(f"PDF 多模态处理失败: {e}")
            result["content_parts"].append(f"[PDF 处理错误: {e}]")

        return result

    def _extract_pdf_text(self, doc) -> str:
        """从 PDF 提取文字"""
        text_parts = []
        for page_num, page in enumerate(doc):
            try:
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"--- 第 {page_num + 1} 页 ---\n{text}")
            except Exception as e:
                logger.warning(f"第 {page_num + 1} 页文字提取失败: {e}")

        return "\n".join(text_parts)

    def _extract_pdf_text_by_page(self, doc) -> Dict[int, str]:
        """按页提取文本，用于图片理解时融合上下文。"""
        page_text = {}
        for page_num, page in enumerate(doc):
            try:
                text = page.get_text() or ""
                page_text[page_num] = text.strip()
            except Exception:
                page_text[page_num] = ""
        return page_text

    def _extract_and_process_pdf_images(self, doc, pdf_path: Path, page_text_map: Dict[int, str]) -> Dict[str, Any]:
        """
        从 PDF 中提取图片并用多模态大模型处理。

        算法：
        1) 收集候选图（含页码、尺寸、页内占比、字节数）
        2) 用视觉签名去重（过滤重复 logo、页眉图）
        3) 重要性评分排序（尺寸 + 页内占比 + 文本稀疏度）
        4) 选 Top-K 并注入页内文本上下文进行理解

        策略：处理所有 >= 5KB 的实质性图片，最多处理 MAX_IMAGES_PER_DOC 张，
        每张 LLM 响应限制在 MAX_TOKENS_PER_IMAGE token 以内。
        """
        content_parts = []
        processed_count = 0
        total_found = 0
        selected_count = 0

        try:
            candidates = self._collect_pdf_image_candidates(doc, page_text_map)
            total_found = len(candidates)
            selected = self._select_pdf_images(candidates)
            selected_count = len(selected)

            logger.info(
                "PDF 图片候选=%s，选中=%s，处理上限=%s",
                total_found,
                selected_count,
                self.MAX_IMAGES_PER_DOC,
            )

            for idx, candidate in enumerate(selected):
                page_num = candidate["page_num"]
                img_index = candidate["img_index"]
                image_bytes = candidate["image_bytes"]
                image_ext = candidate["image_ext"]

                self._emit_progress(
                    f"📄 PDF 图片 [{idx+1}/{selected_count}] "
                    f"(第{page_num+1}页, score={candidate['score']:.2f})"
                )

                page_context = (page_text_map.get(page_num, "") or "")[:self.PAGE_TEXT_CONTEXT_CHARS]
                prompt_context = ""
                if page_context:
                    prompt_context = (
                        "\n\n[页内文本上下文]\n"
                        f"第 {page_num + 1} 页文本片段：\n{page_context}"
                    )

                description = self._process_image_bytes_with_llm(
                    image_bytes,
                    f"PDF第{page_num + 1}页第{img_index + 1}张图片",
                    image_ext,
                    extra_context=prompt_context,
                )

                if description:
                    content_parts.append(
                        f"--- 第 {page_num + 1} 页图片 {img_index + 1} (score={candidate['score']:.2f}) ---\n{description}"
                    )
                    processed_count += 1
                    logger.info(
                        "成功处理第 %s/%s 张选中图片：page=%s img=%s",
                        idx + 1,
                        selected_count,
                        page_num + 1,
                        img_index + 1,
                    )

        except Exception as e:
            logger.error(f"PDF 图片提取失败: {e}")

        logger.info(f"PDF 图片处理完成：共发现 {total_found} 张，处理 {processed_count} 张")
        return {
            "content": "\n".join(content_parts),
            "count": processed_count,
            "found": total_found,
            "selected": selected_count,
        }

    def _collect_pdf_image_candidates(self, doc, page_text_map: Dict[int, str]) -> List[Dict[str, Any]]:
        """收集 PDF 图片候选，并附带用于排序的特征。"""
        candidates: List[Dict[str, Any]] = []
        seen_signatures = set()

        for page_num, page in enumerate(doc):
            try:
                image_list = page.get_images(full=True)
            except Exception as e:
                logger.warning(f"第 {page_num + 1} 页图片列表提取失败: {e}")
                continue

            if not image_list:
                continue

            page_area = max(page.rect.width * page.rect.height, 1)

            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image.get("image", b"")
                    if len(image_bytes) < self.MIN_IMAGE_SIZE:
                        continue

                    image_ext = base_image.get("ext", "png")
                    img_width = img_info[2] if len(img_info) > 2 else 0
                    img_height = img_info[3] if len(img_info) > 3 else 0

                    # 页内占比，越大通常信息密度越高
                    coverage_ratio = 0.0
                    try:
                        rects = page.get_image_rects(xref)
                        if rects:
                            largest = max(rects, key=lambda r: r.width * r.height)
                            coverage_ratio = min((largest.width * largest.height) / page_area, 1.0)
                    except Exception:
                        coverage_ratio = 0.0

                    signature = self._compute_image_signature(image_bytes)
                    if signature in seen_signatures:
                        continue
                    seen_signatures.add(signature)

                    score = self._score_pdf_image_candidate(
                        image_size=len(image_bytes),
                        image_width=img_width,
                        image_height=img_height,
                        coverage_ratio=coverage_ratio,
                        page_text_len=len(page_text_map.get(page_num, "")),
                    )

                    candidates.append(
                        {
                            "page_num": page_num,
                            "img_index": img_index,
                            "image_bytes": image_bytes,
                            "image_ext": image_ext,
                            "score": score,
                            "coverage_ratio": coverage_ratio,
                        }
                    )

                except Exception as e:
                    logger.warning(f"第 {page_num + 1} 页图片 {img_index + 1} 收集失败: {e}")

        return candidates

    def _compute_image_signature(self, image_bytes: bytes) -> str:
        """生成用于去重的图片签名（优先感知哈希，失败则回退字节哈希）。"""
        try:
            from PIL import Image

            img = Image.open(io.BytesIO(image_bytes)).convert("L").resize((16, 16), Image.LANCZOS)
            pixels = list(img.getdata())
            avg = sum(pixels) / max(len(pixels), 1)
            bits = "".join("1" if p > avg else "0" for p in pixels)
            return hashlib.sha256(bits.encode("utf-8")).hexdigest()[:24]
        except Exception:
            return hashlib.sha256(image_bytes[:4096]).hexdigest()[:24]

    def _score_pdf_image_candidate(
        self,
        image_size: int,
        image_width: int,
        image_height: int,
        coverage_ratio: float,
        page_text_len: int,
    ) -> float:
        """给 PDF 图片候选打分，值越大越优先。"""
        size_score = min(image_size / (200 * 1024), 1.0)
        area = max(image_width * image_height, 1)
        area_score = min(area / (900 * 900), 1.0)
        coverage_score = min(max(coverage_ratio, 0.0), 1.0)

        # 文本越少，图片语义权重越高（扫描件、图表页）
        text_sparse_score = 1.0 if page_text_len < 300 else 0.0

        return (
            0.35 * size_score
            + 0.25 * area_score
            + 0.25 * coverage_score
            + 0.15 * text_sparse_score
        )

    def _select_pdf_images(self, candidates: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """按分数选择最有价值的图片，保证覆盖并控制预算。"""
        if not candidates:
            return []

        candidates = sorted(candidates, key=lambda x: x["score"], reverse=True)

        dynamic_limit = max(self.MIN_IMAGE_RETAIN, int(len(candidates) * 0.5))
        limit = min(self.MAX_IMAGES_PER_DOC, dynamic_limit)
        return candidates[:limit]

    def _render_and_process_pdf_pages(self, doc, max_pages: int = 3, page_text_map: Optional[Dict[int, str]] = None) -> str:
        """
        将 PDF 页面渲染为图片并用多模态大模型处理

        Args:
            doc: PyMuPDF 文档对象
            max_pages: 最多处理的页数（避免 token 消耗太大）

        Returns:
            多模态大模型对页面的理解结果
        """
        content_parts = []
        max_pages = min(max_pages, self.MAX_PDF_PAGE_RENDER, len(doc))

        # 优先处理文本稀疏页，可提升扫描件理解准确率。
        page_indices = list(range(len(doc)))
        if page_text_map:
            page_indices.sort(key=lambda i: len(page_text_map.get(i, "")))
        page_indices = page_indices[:max_pages]

        try:
            for page_num in page_indices:
                try:
                    page = doc[page_num]

                    # 渲染页面为图片（2x 缩放以获得更好的清晰度）
                    mat = fitz.Matrix(2, 2)
                    pix = page.get_pixmap(matrix=mat)

                    # 转换为 bytes
                    img_bytes = pix.tobytes("png")

                    # 用多模态大模型处理
                    description = self._process_image_bytes_with_llm(
                        img_bytes,
                        f"PDF第{page_num + 1}页",
                        "png",
                        extra_context=(
                            "\n\n[页内文本上下文]\n"
                            + (page_text_map.get(page_num, "")[:self.PAGE_TEXT_CONTEXT_CHARS] if page_text_map else "")
                        )
                    )

                    if description:
                        content_parts.append(
                            f"--- 第 {page_num + 1} 页渲染内容 ---\n{description}"
                        )
                        logger.info(f"成功处理第 {page_num + 1} 页")

                except Exception as e:
                    logger.warning(f"渲染处理第 {page_num + 1} 页失败: {e}")

        except Exception as e:
            logger.error(f"PDF 页面渲染失败: {e}")

        return "\n".join(content_parts)

    def _process_pptx_multimodal(self, pptx_path: Path) -> Dict[str, Any]:
        """
        处理 PPTX 文件 - 同时提取文字和图片，用多模态大模型理解图片

        Returns:
            包含 content_parts 和 images_processed 的字典
        """
        result = {
            "content_parts": [],
            "images_processed": 0
        }

        try:
            from pptx import Presentation
            from PIL import Image
            import io
        except ImportError as e:
            logger.warning(f"缺少依赖库: {e}，无法进行 PPTX 多模态处理")
            return result

        try:
            prs = Presentation(str(pptx_path))
            logger.info(f"开始处理 PPTX，共 {len(prs.slides)} 张幻灯片")

            # 1. 提取文字内容
            text_content = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text_content.append(f"--- 幻灯片 {slide_num + 1} ---\n" + "\n".join(slide_text))

            if text_content:
                result["content_parts"].append("=== PPTX 文字内容 ===\n" + "\n".join(text_content))
                logger.info(f"提取到 PPTX 文字: {len(text_content)} 张幻灯片")

            # 2. 提取并处理图片
            images_result = self._extract_and_process_pptx_images(prs, pptx_path)
            if images_result["content"].strip():
                result["content_parts"].append("=== PPTX 图片内容 ===\n" + images_result["content"])
            result["images_processed"] = images_result["count"]

        except Exception as e:
            logger.error(f"PPTX 多模态处理失败: {e}")
            result["content_parts"].append(f"[PPTX 处理错误: {e}]")

        return result

    def _extract_and_process_pptx_images(self, prs, pptx_path: Path) -> Dict[str, Any]:
        """
        从 PPTX 中提取图片并用多模态大模型处理（限制数量）

        Returns:
            包含 content 和 count 的字典
        """
        content_parts = []
        processed_count = 0

        try:
            for slide_num, slide in enumerate(prs.slides):
                if processed_count >= self.MAX_IMAGES_PER_DOC:
                    logger.info(f"已达到 PPTX 图片处理上限 ({self.MAX_IMAGES_PER_DOC})，跳过剩余")
                    break

                try:
                    for shape in slide.shapes:
                        if processed_count >= self.MAX_IMAGES_PER_DOC:
                            break

                        # 检查是否有图片
                        if hasattr(shape, "image"):
                            try:
                                image = shape.image
                                image_bytes = image.blob
                                image_ext = image.ext

                                # 跳过太小的图片
                                if len(image_bytes) < self.MIN_IMAGE_SIZE:
                                    continue

                                description = self._process_image_bytes_with_llm(
                                    image_bytes,
                                    f"PPT第{slide_num + 1}页图片",
                                    image_ext
                                )

                                if description:
                                    content_parts.append(
                                        f"--- 第 {slide_num + 1} 页图片 ---\n{description}"
                                    )
                                    processed_count += 1
                                    logger.info(f"成功处理第 {slide_num + 1} 页图片（第 {processed_count} 张）")

                            except Exception as e:
                                logger.warning(f"处理第 {slide_num + 1} 页图片失败: {e}")

                        # 检查是否是图片占位符（PictureFrame）
                        if shape.shape_type == 13:
                            if processed_count >= self.MAX_IMAGES_PER_DOC:
                                break
                            try:
                                image = shape.image
                                image_bytes = image.blob
                                image_ext = image.ext

                                if len(image_bytes) < self.MIN_IMAGE_SIZE:
                                    continue

                                description = self._process_image_bytes_with_llm(
                                    image_bytes,
                                    f"PPT第{slide_num + 1}页PictureFrame图片",
                                    image_ext
                                )

                                if description:
                                    content_parts.append(
                                        f"--- 第 {slide_num + 1} 页 PictureFrame 图片 ---\n{description}"
                                    )
                                    processed_count += 1

                            except Exception as e:
                                logger.warning(f"处理第 {slide_num + 1} 页 PictureFrame 失败: {e}")

                except Exception as e:
                    logger.warning(f"第 {slide_num + 1} 页处理失败: {e}")

        except Exception as e:
            logger.error(f"PPTX 图片提取失败: {e}")

        logger.info(f"PPTX 图片处理完成：共处理 {processed_count} 张")
        return {"content": "\n".join(content_parts), "count": processed_count}

    # ---- 视频处理参数 ----
    VIDEO_MAX_KEYFRAMES_LLM = 12      # 送入 LLM 理解的最大关键帧数
    VIDEO_MAX_TOKENS_PER_FRAME = 384  # 单帧 LLM 响应限制 token
    VIDEO_BATCH_SIZE = 3              # 每批打包发送的帧数（多图单次请求）

    def _process_video_multimodal(self, video_path: Path) -> Dict[str, Any]:
        """
        处理视频文件 - 关键帧提取 + 多模态 LLM 理解

        算法：
        1. 用 VideoProcessor 提取视频元数据和关键帧
        2. 对关键帧压缩后，按批次发送给多模态 LLM 理解
        3. 可选：提取音频并转写（需要 ffmpeg + whisper）
        4. 合成元数据 + 帧描述 + 音频转写

        Returns:
            {content_parts, frames_processed, duration_sec, has_audio_transcript}
        """
        result: Dict[str, Any] = {
            "content_parts": [],
            "frames_processed": 0,
            "duration_sec": None,
            "has_audio_transcript": False,
        }

        try:
            from .video_processor import VideoProcessor
        except ImportError as e:
            logger.error(f"VideoProcessor 导入失败: {e}")
            result["content_parts"].append("[视频处理需要 opencv-python]")
            return result

        vp = VideoProcessor()

        # 1. 提取元数据
        try:
            video_meta = vp._extract_video_metadata(str(video_path))
            result["duration_sec"] = video_meta.get("duration_sec")
            meta_text = vp._format_metadata(video_meta)
            result["content_parts"].append(f"=== 视频信息 ===\n{meta_text}")
            logger.info(f"视频元数据: {meta_text}")
        except Exception as e:
            logger.warning(f"视频元数据提取失败: {e}")
            result["content_parts"].append(f"[视频元数据提取失败: {e}]")

        # 2. 提取关键帧
        try:
            keyframes = vp.extract_keyframes(str(video_path))
            if not keyframes:
                result["content_parts"].append("[未能从视频中提取到关键帧]")
                return result

            # 限制送入 LLM 的帧数
            if len(keyframes) > self.VIDEO_MAX_KEYFRAMES_LLM:
                # 优先保留场景变化帧
                scene_frames = [kf for kf in keyframes if kf["source"] == "scene_change"]
                other_frames = [kf for kf in keyframes if kf["source"] != "scene_change"]

                selected = []
                # 始终保留首帧和尾帧
                if keyframes[0]["source"] == "first":
                    selected.append(keyframes[0])
                if keyframes[-1]["source"] == "last":
                    selected.append(keyframes[-1])

                budget = self.VIDEO_MAX_KEYFRAMES_LLM - len(selected)
                # 优先取场景变化帧
                sc_budget = min(len(scene_frames), int(budget * 0.7))
                step = max(1, len(scene_frames) // sc_budget) if sc_budget > 0 else 1
                selected.extend(scene_frames[::step][:sc_budget])

                # 剩余预算取均匀帧
                remaining = self.VIDEO_MAX_KEYFRAMES_LLM - len(selected)
                if remaining > 0 and other_frames:
                    step = max(1, len(other_frames) // remaining)
                    selected.extend(other_frames[::step][:remaining])

                # 按时间排序
                keyframes = sorted(selected, key=lambda x: x["timestamp_sec"])

            logger.info(f"将对 {len(keyframes)} 个关键帧进行 LLM 理解")

        except Exception as e:
            logger.error(f"关键帧提取失败: {e}")
            result["content_parts"].append(f"[关键帧提取失败: {e}]")
            return result

        # 3. 分批发送关键帧到多模态 LLM 理解
        frame_descriptions = []
        batch_size = self.VIDEO_BATCH_SIZE
        total_batches = (len(keyframes) + batch_size - 1) // batch_size

        for batch_idx in range(total_batches):
            batch_start = batch_idx * batch_size
            batch_end = min(batch_start + batch_size, len(keyframes))
            batch = keyframes[batch_start:batch_end]

            self._emit_progress(
                f"🎬 视频帧批次 [{batch_idx+1}/{total_batches}] "
                f"({len(batch)} 帧)"
            )

            try:
                batch_desc = self._process_video_frame_batch(batch, batch_idx, total_batches)
                if batch_desc:
                    frame_descriptions.append(batch_desc)
                    result["frames_processed"] += len(batch)
                    logger.info(
                        f"视频帧批次 {batch_idx + 1}/{total_batches} 处理完成，"
                        f"帧数={len(batch)}"
                    )
            except Exception as e:
                logger.warning(f"视频帧批次 {batch_idx + 1} 处理失败: {e}")

        if frame_descriptions:
            result["content_parts"].append(
                "=== 视频关键帧内容 ===\n" + "\n\n".join(frame_descriptions)
            )

        # 4. 可选：音频转写
        audio_transcript = self._try_extract_audio_transcript(video_path)
        if audio_transcript:
            result["content_parts"].append(
                "=== 音频转写内容 ===\n" + audio_transcript
            )
            result["has_audio_transcript"] = True

        # 5. 最终综合：如果有多个帧描述，请 LLM 做整体总结
        if result["frames_processed"] >= 3:
            try:
                summary = self._summarize_video_content(
                    frame_descriptions, audio_transcript, video_meta
                )
                if summary:
                    result["content_parts"].insert(1, "=== 视频内容总结 ===\n" + summary)
            except Exception as e:
                logger.warning(f"视频总结生成失败: {e}")

        return result

    def _process_video_frame_batch(
        self, batch: List[Dict[str, Any]], batch_idx: int, total_batches: int
    ) -> str:
        """
        处理一批视频关键帧，发送给多模态 LLM

        多帧打包发送可以让 LLM 理解帧之间的时间关系。
        """
        if not batch:
            return ""

        # 构建包含多张图片的 content 列表
        content_items = []

        # 文本 prompt
        timestamps = ", ".join(
            f"{kf['timestamp_str']}({kf['source']})" for kf in batch
        )
        prompt = (
            f"以下是视频的第 {batch_idx + 1}/{total_batches} 组关键帧，"
            f"时间戳: [{timestamps}]。\n\n"
            f"请分析每个关键帧的内容：\n"
            f"1. 描述画面中的主要内容（人物、场景、动作、文字等）\n"
            f"2. 如果画面中有文字或字幕，请完整转录\n"
            f"3. 识别画面类型（实景、动画、PPT、屏幕录制、字幕卡等）\n"
            f"4. 如果能看出情节或主题变化，请说明\n\n"
            f"请按时间顺序逐帧描述，用中文回答。"
        )
        content_items.append({"type": "text", "text": prompt})

        # 逐张添加图片
        for kf in batch:
            compressed = self._compress_image_for_llm(kf["frame_bytes"])
            b64 = base64.b64encode(compressed).decode("utf-8")
            content_items.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{b64}"},
            })

        # 调用 LLM
        if not self.client:
            return "[多模态处理不可用]"

        try:
            max_tokens = self.VIDEO_MAX_TOKENS_PER_FRAME * len(batch)

            def do_call():
                return self.client.chat.completions.create(
                    model=self.model,
                    messages=[{"role": "user", "content": content_items}],
                    max_tokens=max_tokens,
                )

            response, timed_out = self._call_with_timeout(
                do_call, timeout=self.LLM_TIMEOUT
            )
            if timed_out:
                logger.warning(f"视频帧批次 LLM 调用超时（{self.LLM_TIMEOUT}s）")
                return "[视频帧处理超时]"
            if response is None:
                return "[视频帧处理失败]"

            description = response.choices[0].message.content
            # 为输出加上时间戳标注
            header = " | ".join(
                f"[{kf['timestamp_str']}]" for kf in batch
            )
            return f"--- 帧组 {batch_idx + 1}: {header} ---\n{description}"

        except Exception as e:
            logger.error(f"视频帧 LLM 处理失败: {e}")
            return f"[视频帧处理失败: {e}]"

    def _try_extract_audio_transcript(self, video_path: Path) -> Optional[str]:
        """
        尝试提取视频音频并转写

        需要 ffmpeg + whisper。如果不可用，静默跳过。
        """
        try:
            from .video_processor import VideoProcessor
            audio_path = VideoProcessor.extract_audio_track(str(video_path))
            if not audio_path:
                logger.info("无法提取音频轨道（可能缺少 ffmpeg），跳过音频转写")
                return None

            # 尝试用 whisper 转写
            try:
                import whisper
                model = whisper.load_model("base")
                result = model.transcribe(audio_path, language="zh")
                text = result.get("text", "").strip()
                if text:
                    logger.info(f"音频转写成功：{len(text)} 字符")
                    return text
            except ImportError:
                logger.info("whisper 未安装，跳过音频转写")
            except Exception as e:
                logger.warning(f"whisper 转写失败: {e}")

            # 清理临时音频文件
            try:
                import os
                os.remove(audio_path)
            except OSError:
                pass

        except Exception as e:
            logger.debug(f"音频转写尝试失败: {e}")

        return None

    def _summarize_video_content(
        self,
        frame_descriptions: List[str],
        audio_transcript: Optional[str],
        video_meta: Dict[str, Any],
    ) -> str:
        """使用文本 LLM 综合视频帧描述和音频转写，生成整体总结"""
        parts = []
        parts.append(f"视频时长: {video_meta.get('duration_str', '未知')}")

        if frame_descriptions:
            combined_frames = "\n".join(frame_descriptions)
            # 限制长度避免超 token
            if len(combined_frames) > 8000:
                combined_frames = combined_frames[:8000] + "\n[...截断]"
            parts.append(f"关键帧描述:\n{combined_frames}")

        if audio_transcript:
            transcript = audio_transcript[:4000]
            parts.append(f"音频转写:\n{transcript}")

        prompt = (
            "请根据以下视频关键帧描述" +
            ("和音频转写内容" if audio_transcript else "") +
            "，生成一段完整的视频内容总结。\n\n"
            "要求：\n"
            "1. 概述视频的主题和主要内容\n"
            "2. 按时间线梳理关键信息\n"
            "3. 如果是教学/演讲视频，提取核心观点\n"
            "4. 如果有屏幕录制或PPT内容，提取关键信息\n"
            "5. 用中文回答，300-500字\n\n"
            + "\n".join(parts)
        )

        return self._call_text_llm(prompt)

    def _process_image_with_llm(self, image_path: Path) -> str:
        """使用多模态大模型直接理解图片内容（两阶段：感知→分析）"""
        image_bytes = image_path.read_bytes()
        image_ext = image_path.suffix.lstrip('.')
        self._emit_progress(f"🖼️  正在理解图片: {image_path.name}")

        return self._process_image_bytes_with_llm(
            image_bytes,
            image_path.name,
            image_ext
        )

    def _compress_image_for_llm(self, image_bytes: bytes) -> bytes:
        """
        Claude Code 风格的智能图片压缩算法（优化版 v2）

        核心策略：
        1. 小图直接返回（避免无谓压缩 I/O）
        2. 单次读取，多次复用
        3. 渐进质量：先尝试高质量，若太大再降级
        4. Resize 到合理尺寸（600px），足够内容理解

        参数可通过类常量配置：
        - COMPRESS_THRESHOLD_KB: 触发压缩的最小图片大小
        - COMPRESS_MAX_DIM: resize 最大尺寸
        - COMPRESS_TARGET_KB: 压缩目标大小
        - COMPRESS_QUALITIES: 渐进质量列表
        """
        try:
            from PIL import Image
            import io
        except ImportError:
            return image_bytes  # 无法压缩，原样返回

        try:
            # 阶段1：检查原始大小 - 小于阈值直接返回（避免无谓处理）
            threshold = getattr(self, 'COMPRESS_THRESHOLD_KB', 50) * 1024
            if len(image_bytes) < threshold:
                return image_bytes

            # 单次读取，复用 buffer
            img = Image.open(io.BytesIO(image_bytes))
            original_w, original_h = img.size

            # 阶段2：检查原始尺寸 - 已经很小的图片直接压缩返回
            max_dim = getattr(self, 'COMPRESS_MAX_DIM', 600)
            if original_w <= max_dim and original_h <= max_dim:
                buf = io.BytesIO()
                img.save(buf, format='JPEG', quality=70)
                compressed = buf.getvalue()
                # 如果压缩后反而更大，返回原图
                return compressed if len(compressed) < len(image_bytes) else image_bytes

            # 阶段3：Resize 到 max_dim（保持宽高比）
            # 使用 LANCZOS 重采样（高质量缩放）
            if original_w > max_dim or original_h > max_dim:
                ratio = min(max_dim / original_w, max_dim / original_h)
                new_w = int(original_w * ratio)
                new_h = int(original_h * ratio)
                img = img.resize((new_w, new_h), Image.LANCZOS)

            # 阶段4：渐进质量压缩
            qualities = getattr(self, 'COMPRESS_QUALITIES', [70, 50, 30])
            target_kb = getattr(self, 'COMPRESS_TARGET_KB', 80) * 1024
            for quality in qualities:
                buf = io.BytesIO()
                img.save(buf, format='JPEG', quality=quality, optimize=True)
                compressed = buf.getvalue()
                # 目标：压到目标大小以下
                if len(compressed) < target_kb:
                    return compressed
                # 如果已经是最低质量或比原图还大，返回原图
                if quality == qualities[-1]:
                    return compressed if len(compressed) < len(image_bytes) else image_bytes

            return image_bytes

        except Exception as e:
            logger.warning(f"图片压缩失败: {e}")
            return image_bytes  # 压缩失败，返回原图

    def _process_image_bytes_with_llm(
        self,
        image_bytes: bytes,
        image_name: str,
        image_ext: str,
        extra_context: str = "",
    ) -> str:
        """
        两阶段图片理解流程：
          阶段1 感知层 — 多模态LLM识别图片类型和内容
          阶段2 分析层 — 根据图片类型用针对性prompt提取实体/概念

        合并为单次 LLM 调用以节省延迟（prompt 内嵌类型判断 + 对应提取策略）。
        """
        compressed = self._compress_image_for_llm(image_bytes)
        if compressed != image_bytes:
            logger.debug(f"图片压缩: {len(image_bytes)/1024:.0f}KB → {len(compressed)/1024:.0f}KB")
        image_bytes = compressed

        base64_image = base64.b64encode(image_bytes).decode('utf-8')

        prompt = f"""请分析这张图片（{image_name}），按以下两个阶段完成：

## 阶段1：感知识别
首先判断图片类型（从以下选项中选择一项）：
- 图表(chart) — 柱状图、折线图、饼图、散点图等数据可视化
- 表格(table) — 结构化数据表格
- 示意图(diagram) — 流程图、架构图、UML、思维导图等
- 文档扫描(scan) — 扫描件、截图、包含大段文字的图片
- 照片(photo) — 实景照片（人物、场景、物体）
- 截图(screenshot) — UI界面、网页、代码截图
- 公式(formula) — 数学公式、化学方程式
- 其他(other)

然后给出图片的简要描述（1-2句话）。

## 阶段2：针对性内容提取
根据阶段1判断的类型，执行对应的提取策略：

**如果是图表(chart)**：提取图表标题、坐标轴含义、数据趋势、关键数值、数据对比结论。
**如果是表格(table)**：用 Markdown 格式完整呈现表格数据，提取表头含义。
**如果是示意图(diagram)**：识别所有节点/模块名称、它们之间的连接关系和方向、整体架构含义。
**如果是文档扫描(scan)**：完整转录所有可见文字，保留段落结构。
**如果是照片(photo)**：识别人物（姓名/身份如可判断）、场景、物体、事件、地点。
**如果是截图(screenshot)**：识别应用/网站名称、UI元素、显示的数据/文字。
**如果是公式(formula)**：用 LaTeX 格式转录公式，解释其含义和变量。
**如果是其他(other)**：描述主要内容，提取所有可见文字。

最后，列出图片中涉及的**实体**（人名、公司、产品、地名等专有名词）和**概念**（方法、理论、技术术语等）。

请用中文回答，确保信息完整准确。如果图片质量较低或无法识别，请明确说明。"""

        if extra_context:
            prompt += f"\n\n上下文参考：{extra_context}"

        return self._call_multimodal_llm(
            prompt=prompt,
            base64_image=base64_image,
            image_ext=image_ext
        )

    def _process_other_with_llm(self, path: Path, ext: str) -> str:
        """处理其他类型的文件"""
        # 先尝试读取文件内容
        try:
            content = path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            content = path.read_text(encoding='gbk', errors='ignore')
        except Exception as e:
            return f"文件读取失败: {e}"

        # 如果内容太长，进行预处理
        if len(content) > 50000:
            content = content[:50000] + "\n[内容过长，已截断]"

        prompt = f"""请分析以下{ext}文件的内容，并提供：
1. 文档结构总结
2. 关键信息提取
3. 如果包含表格，请用 Markdown 格式呈现

文件内容：
{content}"""

        return self._call_text_llm(prompt)

    def _encode_image_to_base64(self, image_path: Path) -> str:
        """将图片编码为 base64 字符串"""
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')

    def _call_with_timeout(self, func, *args, timeout: int = None, **kwargs) -> Any:
        """使用线程在指定秒数内执行函数，超时则中断"""
        timeout_val = timeout or self.LLM_TIMEOUT
        result = {"value": None, "error": None}

        def target():
            try:
                result["value"] = func(*args, **kwargs)
            except Exception as e:
                result["error"] = e

        thread = threading.Thread(target=target)
        thread.daemon = True
        thread.start()
        thread.join(timeout_val)

        if thread.is_alive():
            # 超时，线程还在运行
            return None, True  # (value, timed_out)
        if result["error"]:
            return None, False  # (value, timed_out)
        return result["value"], False  # (value, timed_out)

    def _call_multimodal_llm(self, prompt: str, base64_image: str, image_ext: str = "png") -> str:
        """调用多模态大模型（支持图片），带超时控制"""
        if not self.client:
            return "[多模态处理不可用，需要配置 LLM]"

        try:
            content = [
                {"type": "text", "text": prompt},
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/{image_ext};base64,{base64_image}"
                    }
                }
            ]

            def do_call():
                return self.client.chat.completions.create(
                    model=self.model,
                    messages=[{"role": "user", "content": content}],
                    max_tokens=self.MAX_TOKENS_PER_IMAGE,
                )

            response, timed_out = self._call_with_timeout(do_call, timeout=self.LLM_TIMEOUT)
            if timed_out:
                logger.warning(f"LLM 调用超时（{self.LLM_TIMEOUT}秒），跳过此图片")
                return "[图片处理超时，已跳过]"
            if response is None:
                return "[多模态处理失败]"

            return response.choices[0].message.content

        except Exception as e:
            logger.error(f"多模态 LLM 调用失败: {e}")
            return f"[多模态处理失败: {e}]"

    def _call_text_llm(self, prompt: str) -> str:
        """调用文本 LLM（不使用图片）"""
        if not self.client:
            return "[LLM 处理不可用]"

        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=4096,
            )

            return response.choices[0].message.content

        except Exception as e:
            logger.error(f"文本 LLM 调用失败: {e}")
            return f"[LLM 处理失败: {e}]"

    def _process_fallback(self, path: Path, ext: str) -> ProcessedDocument:
        """回退到简单处理方法"""
        content = ""
        errors = []
        metadata = {
            "source": str(path),
            "type": ext.lstrip('.'),
            "processor": "fallback",
        }

        try:
            # 尝试作为文本读取
            content = path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            try:
                content = path.read_text(encoding='gbk', errors='ignore')
            except Exception as e:
                errors.append(f"文本读取失败: {e}")
                content = f"[{ext.upper()} 文件，无法直接读取]"
        except Exception as e:
            errors.append(f"文件读取失败: {e}")
            content = f"[{ext.upper()} 文件]"

        return ProcessedDocument(
            content=content,
            metadata=metadata,
            errors=errors
        )

    def enhance_content_with_llm(self, content: str, file_type: str) -> str:
        """使用 LLM 增强已提取的内容"""
        if not self.client:
            return content

        prompt = f"""请优化以下从{file_type}文件中提取的内容，使其更易读、结构更清晰：

1. 修复格式问题
2. 整理段落和标题
3. 如果有表格数据，确保格式正确
4. 保留所有原始信息

原始内容：
{content}"""

        return self._call_text_llm(prompt)
