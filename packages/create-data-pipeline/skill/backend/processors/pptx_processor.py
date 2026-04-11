"""
PowerPoint PPTX 文件处理器 - 提取演示文稿内容
"""

import logging
from pathlib import Path
from typing import List, Dict, Any
from .file_processor import ProcessedDocument

logger = logging.getLogger(__name__)


class PptxProcessor:
    """PowerPoint PPTX 演示文稿处理器"""

    def __init__(self):
        self.extract_images = False

    def process(self, file_path: str) -> ProcessedDocument:
        """处理 PowerPoint PPTX 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            return ProcessedDocument(
                content="",
                errors=["缺少依赖: python-pptx。请运行: pip install python-pptx"]
            )

        path = Path(file_path)
        errors = []
        content_parts = []

        try:
            prs = Presentation(str(path))

            # 提取元数据
            metadata = {
                "source": str(path),
                "type": "pptx",
                "slides": len(prs.slides),
                "title": prs.core_properties.title or "",
            }

            # 提取每一页的内容
            for slide_num, slide in enumerate(prs.slides):
                slide_content = []

                # 提取标题（通常是第一个文本框）
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                if slide_content:
                    content_parts.append(f"\\n--- 幻灯片 {slide_num + 1} ---\\n")
                    content_parts.append("\\n".join(slide_content))

            content = "\\n".join(content_parts)

            return ProcessedDocument(
                content=content,
                metadata=metadata,
                errors=errors
            )

        except Exception as e:
            logger.error(f"PPTX 处理失败: {e}")
            return ProcessedDocument(
                content="",
                errors=[f"PPTX 处理失败: {str(e)}"]
            )
"""
PowerPoint PPTX 文件处理器 - 提取演示文稿内容
"""

import logging
from pathlib import Path
from typing import List, Dict, Any
from .file_processor import ProcessedDocument

logger = logging.getLogger(__name__)


class PptxProcessor:
    """PowerPoint PPTX 演示文稿处理器"""

    def __init__(self):
        self.extract_images = False

    def process(self, file_path: str) -> ProcessedDocument:
        """处理 PowerPoint PPTX 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            return ProcessedDocument(
                content="",
                errors=["缺少依赖: python-pptx。请运行: pip install python-pptx"]
            )

        path = Path(file_path)
        errors = []
        content_parts = []

        try:
            prs = Presentation(str(path))

            # 提取元数据
            metadata = {
                "source": str(path),
                "type": "pptx",
                "slides": len(prs.slides),
                "title": prs.core_properties.title or "",
            }

            # 提取每一页的内容
            for slide_num, slide in enumerate(prs.slides):
                slide_content = []

                # 提取标题（通常是第一个文本框）
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                if slide_content:
                    content_parts.append(f"\\n--- 幻灯片 {slide_num + 1} ---\\n")
                    content_parts.append("\\n".join(slide_content))

            content = "\\n".join(content_parts)

            return ProcessedDocument(
                content=content,
                metadata=metadata,
                errors=errors
            )

        except Exception as e:
            logger.error(f"PPTX 处理失败: {e}")
            return ProcessedDocument(
                content="",
                errors=[f"PPTX 处理失败: {str(e)}"]
            )
