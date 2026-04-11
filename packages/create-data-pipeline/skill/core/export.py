"""
数据导出模块 - 将 wiki 数据资产导出为 PPT / Excel / Markdown 等格式

支持用户自定义意图的数据消费：
- 从 wiki 页面导出为结构化 Excel
- 从 wiki 合成内容导出为 PPT 演示文稿
- 从 graph 数据导出关系矩阵
"""

import json
import re
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import date


class WikiExporter:
    """wiki 数据导出器 - 将 wiki 数据资产转化为可消费的格式"""

    def __init__(self, wiki_dir: Path, graph_dir: Optional[Path] = None):
        self.wiki_dir = wiki_dir
        self.graph_dir = graph_dir or wiki_dir.parent / "graph"

    def collect_pages(self, page_type: Optional[str] = None) -> List[Dict[str, Any]]:
        """
        收集 wiki 页面信息

        Args:
            page_type: 过滤类型 (source/entity/concept/synthesis)，None=全部

        Returns:
            [{path, title, type, tags, content, sources, last_updated}, ...]
        """
        pages = []
        for p in self.wiki_dir.rglob("*.md"):
            if p.name in ("index.md", "log.md", "lint-report.md"):
                continue
            try:
                content = p.read_text(encoding="utf-8")
                meta = self._parse_frontmatter(content)
                if page_type and meta.get("type") != page_type:
                    continue
                pages.append({
                    "path": str(p.relative_to(self.wiki_dir)),
                    "title": meta.get("title", p.stem),
                    "type": meta.get("type", "unknown"),
                    "tags": meta.get("tags", []),
                    "sources": meta.get("sources", []),
                    "last_updated": meta.get("last_updated", ""),
                    "content": content,
                })
            except Exception:
                pass
        return pages

    def export_to_excel(self, output_path: str, page_type: Optional[str] = None,
                        include_content: bool = False) -> str:
        """
        导出 wiki 数据为 Excel

        Args:
            output_path: 输出文件路径
            page_type: 过滤页面类型
            include_content: 是否包含完整内容列

        Returns:
            输出文件路径
        """
        try:
            import openpyxl
            from openpyxl.styles import Font, Alignment, PatternFill
        except ImportError:
            raise ImportError("需要 openpyxl: pip install openpyxl")

        pages = self.collect_pages(page_type)
        wb = openpyxl.Workbook()

        # 概览 sheet
        ws = wb.active
        ws.title = "Wiki 概览"
        headers = ["标题", "类型", "标签", "来源", "最后更新", "文件路径"]
        if include_content:
            headers.append("内容")

        # 表头样式
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)

        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=h)
            cell.fill = header_fill
            cell.font = header_font

        for row_idx, page in enumerate(pages, 2):
            ws.cell(row=row_idx, column=1, value=page["title"])
            ws.cell(row=row_idx, column=2, value=page["type"])
            ws.cell(row=row_idx, column=3, value=", ".join(page["tags"]) if page["tags"] else "")
            ws.cell(row=row_idx, column=4, value=", ".join(page["sources"]) if page["sources"] else "")
            ws.cell(row=row_idx, column=5, value=page["last_updated"])
            ws.cell(row=row_idx, column=6, value=page["path"])
            if include_content:
                # 去掉 frontmatter
                body = self._strip_frontmatter(page["content"])
                ws.cell(row=row_idx, column=7, value=body[:32000])  # Excel 单元格限制

        # 自动列宽
        for col in ws.columns:
            max_len = max(len(str(cell.value or "")) for cell in col)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 60)

        # 关系矩阵 sheet（如果有 graph 数据）
        graph_json = self.graph_dir / "graph.json"
        if graph_json.exists():
            self._add_graph_sheet(wb, graph_json)

        wb.save(output_path)
        return output_path

    def export_to_pptx(self, output_path: str, title: str = "Wiki 知识报告",
                       page_type: Optional[str] = None,
                       custom_query: Optional[str] = None,
                       synthesis_content: Optional[str] = None) -> str:
        """
        导出为 PPT 演示文稿

        Args:
            output_path: 输出文件路径
            title: PPT 标题
            page_type: 过滤页面类型
            custom_query: 自定义汇报主题（用于生成 PPT 标题页）
            synthesis_content: 合成内容（如 query 结果，用于正文）

        Returns:
            输出文件路径
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("需要 python-pptx: pip install python-pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 标题页
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = (
            f"生成日期: {date.today().isoformat()}\n"
            f"数据来源: LLM Wiki Knowledge Base"
        )

        # 如果有合成内容，按 ## 分页
        if synthesis_content:
            sections = re.split(r'\n##\s+', synthesis_content)
            for section in sections:
                if not section.strip():
                    continue
                lines = section.strip().split('\n')
                section_title = lines[0].strip('#').strip()
                section_body = '\n'.join(lines[1:]).strip()

                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = section_title
                body = slide.placeholders[1]
                body.text = section_body[:2000]  # PPT 文本限制

        # 按页面类型生成概览
        pages = self.collect_pages(page_type)
        if pages:
            # 统计页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "知识库统计"
            type_counts = {}
            for p in pages:
                t = p["type"]
                type_counts[t] = type_counts.get(t, 0) + 1
            stats_text = "\n".join(f"• {t}: {c} 个页面" for t, c in sorted(type_counts.items()))
            stats_text += f"\n\n总计: {len(pages)} 个页面"
            slide.placeholders[1].text = stats_text

            # 关键页面
            for page in pages[:10]:  # 最多 10 页
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = page["title"]
                body_text = self._strip_frontmatter(page["content"])
                # 提取前 1500 字符
                slide.placeholders[1].text = body_text[:1500]

        prs.save(output_path)
        return output_path

    def export_graph_data(self) -> Optional[Dict[str, Any]]:
        """导出知识图谱数据（JSON 格式）"""
        graph_json = self.graph_dir / "graph.json"
        if not graph_json.exists():
            return None
        try:
            return json.loads(graph_json.read_text(encoding="utf-8"))
        except Exception:
            return None

    def _parse_frontmatter(self, content: str) -> Dict[str, Any]:
        """解析 YAML frontmatter"""
        meta = {}
        m = re.match(r'^---\s*\n(.*?)\n---', content, re.DOTALL)
        if not m:
            return meta
        fm = m.group(1)
        for line in fm.split('\n'):
            match = re.match(r'^(\w[\w_]*)\s*:\s*(.+)', line)
            if match:
                key = match.group(1)
                val = match.group(2).strip().strip('"').strip("'")
                if val.startswith('[') and val.endswith(']'):
                    val = [v.strip().strip('"').strip("'") for v in val[1:-1].split(',') if v.strip()]
                meta[key] = val
        return meta

    def _strip_frontmatter(self, content: str) -> str:
        """移除 frontmatter，返回正文"""
        return re.sub(r'^---\s*\n.*?\n---\s*\n?', '', content, flags=re.DOTALL).strip()

    def _add_graph_sheet(self, wb, graph_json_path: Path):
        """在 Excel 中添加知识图谱关系矩阵 sheet"""
        try:
            data = json.loads(graph_json_path.read_text(encoding="utf-8"))
            edges = data.get("edges", [])
            if not edges:
                return

            ws = wb.create_sheet("知识关系")
            headers = ["源页面", "目标页面", "关系类型", "关系描述", "置信度"]
            for col, h in enumerate(headers, 1):
                ws.cell(row=1, column=col, value=h)

            for row_idx, edge in enumerate(edges, 2):
                ws.cell(row=row_idx, column=1, value=edge.get("from", ""))
                ws.cell(row=row_idx, column=2, value=edge.get("to", ""))
                ws.cell(row=row_idx, column=3, value=edge.get("type_cn", edge.get("type", "")))
                ws.cell(row=row_idx, column=4, value=edge.get("label", ""))
                ws.cell(row=row_idx, column=5, value=edge.get("confidence", ""))
        except Exception:
            pass
