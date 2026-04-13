# tests/unit/test_inductor.py
import pytest
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from backend.ppt.inductor import SlideInducter

@pytest.fixture
def sample_template(tmp_path):
    """创建测试用 PPT 模板"""
    prs = Presentation()

    # 幻灯片 1: 标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    title.text = "Title Slide"

    # 幻灯片 2: 双栏文本
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    left_box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(4))
    left_box.text = "Left Column"
    right_box = slide2.shapes.add_textbox(Inches(5), Inches(2), Inches(3), Inches(4))
    right_box.text = "Right Column"

    # 保存
    template_path = tmp_path / "test_template.pptx"
    prs.save(str(template_path))
    return str(template_path)

def test_analyze_template(sample_template):
    """测试模板分析"""
    inductor = SlideInducter()
    layouts = inductor.analyze(sample_template)

    # 应该提取出至少 1 种布局
    assert len(layouts) >= 1

    # 检查布局属性
    for layout in layouts:
        assert layout.name is not None
        assert layout.description is not None
        assert isinstance(layout.elements, list)

def test_nonexistent_file():
    """测试不存在的文件"""
    inductor = SlideInducter()

    with pytest.raises(FileNotFoundError):
        inductor.analyze('/nonexistent/file.pptx')

def test_empty_ppt(tmp_path):
    """测试空 PPT（没有形状）"""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])

    template_path = tmp_path / "empty_template.pptx"
    prs.save(str(template_path))

    inductor = SlideInducter()
    layouts = inductor.analyze(str(template_path))

    # 应该提取出布局，即使元素为空
    assert len(layouts) >= 1

def test_mixed_content(tmp_path):
    """测试混合内容（文本和图片）"""
    prs = Presentation()

    # 创建带文本和图片的幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 添加文本框
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = "This is a text box"

    # 添加另一个文本框（模拟标题）
    title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Title Here"
    # 设置大字体
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(36)

    template_path = tmp_path / "mixed_content.pptx"
    prs.save(str(template_path))

    inductor = SlideInducter()
    layouts = inductor.analyze(str(template_path))

    # 应该提取出布局
    assert len(layouts) >= 1

    # 检查是否识别出不同类型的元素
    layout = layouts[0]
    element_types = [elem.type for elem in layout.elements]
    assert 'title' in element_types
